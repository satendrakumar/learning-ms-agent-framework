"""Build the training deck: Microsoft Agent Framework — Building Production-Ready AI Agents.

Generates `slides/Microsoft-Agent-Framework-Training.pptx` from scratch, so the
deck is reproducible and reviewable as code.

Run:  uv run python slides/build_deck.py
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------------------
# Theme
# --------------------------------------------------------------------------------------

BG = RGBColor(0x0B, 0x0F, 0x19)  # deep navy-black canvas
PANEL = RGBColor(0x15, 0x1B, 0x2B)  # raised surface
PANEL_2 = RGBColor(0x1D, 0x25, 0x3A)  # secondary surface
CODE_BG = RGBColor(0x08, 0x0C, 0x14)
STROKE = RGBColor(0x2A, 0x35, 0x4D)

ACCENT = RGBColor(0x00, 0x78, 0xD4)  # Microsoft blue
CYAN = RGBColor(0x50, 0xE6, 0xFF)  # Azure cyan
VIOLET = RGBColor(0x9B, 0x7C, 0xF6)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
ROSE = RGBColor(0xFB, 0x71, 0x85)

TEXT = RGBColor(0xE8, 0xED, 0xF5)
MUTED = RGBColor(0x8B, 0x97, 0xAC)
DIM = RGBColor(0x5A, 0x67, 0x80)

UI = "Segoe UI"
UI_LIGHT = "Segoe UI Light"
MONO = "Consolas"

SW = 13.333
SH = 7.5
ML = 0.78  # left margin
CW = SW - 2 * ML  # content width

# --------------------------------------------------------------------------------------
# Syntax highlighting (lightweight, line-oriented)
# --------------------------------------------------------------------------------------

KW = (
    "async|await|def|class|from|import|return|if|elif|else|for|while|in|not|and|or|is|"
    "None|True|False|with|as|try|except|finally|yield|lambda|pass|raise|global|del"
)

TOKEN_RE = re.compile(
    r"(?P<comment>#.*$)"
    r"|(?P<string>\"\"\".*?\"\"\"|'''.*?'''|\"[^\"]*\"|'[^']*')"
    r"|(?P<deco>@[A-Za-z_][\w.]*)"
    r"|(?P<kw>\b(?:" + KW + r")\b)"
    r"|(?P<cls>\b[A-Z][A-Za-z0-9_]*\b)"
    r"|(?P<num>\b\d[\d_]*\b)"
    r"|(?P<fn>\b[a-z_][a-z0-9_]*(?=\())"
)

SYNTAX = {
    "comment": RGBColor(0x64, 0x74, 0x8B),
    "string": RGBColor(0xE5, 0xA5, 0x7A),
    "deco": RGBColor(0xFF, 0xCE, 0x5C),
    "kw": RGBColor(0x7A, 0xA2, 0xF7),
    "cls": RGBColor(0x50, 0xE6, 0xFF),
    "num": RGBColor(0x86, 0xE1, 0xA0),
    "fn": RGBColor(0xB4, 0xC9, 0xFF),
}
CODE_FG = RGBColor(0xD6, 0xDE, 0xEB)


def tokenize(line: str):
    """Yield (text, color) runs for one line of Python."""
    out, pos = [], 0
    for m in TOKEN_RE.finditer(line):
        if m.start() > pos:
            out.append((line[pos : m.start()], CODE_FG))
        out.append((m.group(0), SYNTAX[m.lastgroup]))
        pos = m.end()
    if pos < len(line):
        out.append((line[pos:], CODE_FG))
    return out or [(line or " ", CODE_FG)]


# --------------------------------------------------------------------------------------
# Deck helpers
# --------------------------------------------------------------------------------------


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(SW)
        self.prs.slide_height = Inches(SH)
        self.module = ""
        self.count = 0
        self.panels: list[tuple[float, float, float, float]] = []  # code panels, this slide
        self.warnings: list[str] = []

    def _check_clear(self, x, y, w, h, what):
        """Warn if a panel-level shape lands on top of a code panel."""
        for px, py, pw, ph in self.panels:
            if x < px + pw and px < x + w and y < py + ph and py < y + h:
                self.warnings.append(
                    f"slide {self.count}: {what} at y={y:.2f} overlaps code panel "
                    f"ending at y={py + ph:.2f}"
                )

    # -- canvas ------------------------------------------------------------------------

    def _blank(self):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = BG
        self.count += 1
        self.panels = []
        return s

    def slide(self, title: str, eyebrow: str | None = None, footer: bool = True):
        """Standard content slide with title, accent rule and footer."""
        s = self._blank()
        eb = eyebrow if eyebrow is not None else self.module
        top = 0.46
        if eb:
            self.text(s, ML, 0.42, CW, 0.28, eb.upper(), 11, CYAN, bold=True)
            top = 0.72
        self.text(s, ML, top, CW, 0.62, title, 30, TEXT, bold=True, font=UI_LIGHT)
        self.rect(s, ML, top + 0.72, 1.15, 0.055, ACCENT)
        if footer:
            self._footer(s)
        return s

    def _footer(self, s):
        self.rect(s, 0, SH - 0.035, SW, 0.035, PANEL_2)
        self.text(
            s, ML, SH - 0.52, 6.0, 0.3,
            "Microsoft Agent Framework  ·  Production-Ready AI Agents",
            9, DIM,
        )
        self.text(
            s, SW - ML - 1.2, SH - 0.52, 1.2, 0.3, f"{self.count:02d}", 9, DIM,
            align=PP_ALIGN.RIGHT,
        )

    # -- primitives --------------------------------------------------------------------

    def text(self, s, x, y, w, h, txt, size, color, *, bold=False, font=UI,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, italic=False):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = txt
        f = r.font
        f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
        f.color.rgb = color
        return tb

    def rect(self, s, x, y, w, h, fill, *, line=None, lw=1.0, radius=None,
             shape=MSO_SHAPE.RECTANGLE):
        sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.shadow.inherit = False
        if fill is None:
            sh.fill.background()
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line
            sh.line.width = Pt(lw)
        if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
            sh.adjustments[0] = radius
        sh.text_frame.word_wrap = True
        return sh

    def card(self, s, x, y, w, h, *, fill=PANEL, line=STROKE, radius=0.06):
        self._check_clear(x, y, w, h, "card")
        return self.rect(s, x, y, w, h, fill, line=line, radius=radius,
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)

    def label(self, sh, lines, *, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
              pad=0.06):
        """Fill a shape with (text, size, color, bold) lines."""
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Inches(pad)
        tf.margin_top = tf.margin_bottom = Inches(0.04)
        for i, (txt, size, color, bold) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = 1.05
            p.space_after = Pt(2)
            r = p.add_run()
            r.text = txt
            r.font.size, r.font.bold, r.font.name = Pt(size), bold, UI
            r.font.color.rgb = color
        return sh

    def arrow(self, s, x, y, w, h=0.26, color=ACCENT, shape=MSO_SHAPE.RIGHT_ARROW):
        a = self.rect(s, x, y, w, h, color, shape=shape)
        a.adjustments[0] = 0.42
        a.adjustments[1] = 0.55
        return a

    # -- composites --------------------------------------------------------------------

    def bullets(self, s, x, y, w, items, *, size=16, gap=0.72, bullet=ACCENT):
        """items: (headline, subtext) — headline bold, subtext muted."""
        cy = y
        for head, sub in items:
            self.rect(s, x, cy + 0.11, 0.09, 0.09, bullet,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            self.text(s, x + 0.28, cy, w - 0.28, 0.3, head, size, TEXT, bold=True)
            if sub:
                self.text(s, x + 0.28, cy + 0.31, w - 0.28, 0.5, sub, size - 3.5, MUTED,
                          spacing=1.1)
            cy += gap if sub else gap - 0.26
        return cy

    def code(self, s, x, y, w, lines, *, size=13.5, title=None, h=None):
        """Dark code panel with syntax highlighting."""
        lh = size * 1.42 / 72
        body_h = h if h else len(lines) * lh + 0.42
        top = y
        if title:
            bar = self.rect(s, x, y, w, 0.34, PANEL_2, radius=0.12,
                            shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            self.label(bar, [(title, 10.5, MUTED, True)], align=PP_ALIGN.LEFT, pad=0.16)
            top = y + 0.30
            body_h += 0.06
        self.rect(s, x, top, w, body_h, CODE_BG, line=STROKE, radius=0.03,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        self.panels.append((x, top, w, body_h))
        self.rect(s, x, top, 0.045, body_h, ACCENT)
        tb = s.shapes.add_textbox(Inches(x + 0.26), Inches(top + 0.18),
                                  Inches(w - 0.4), Inches(body_h - 0.3))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # exact spacing, so the panel height maths above is not an estimate
            p.line_spacing = Pt(size * 1.42)
            p.space_after = Pt(0)
            for txt, color in tokenize(line):
                r = p.add_run()
                r.text = txt
                r.font.size, r.font.name = Pt(size), MONO
                r.font.color.rgb = color
        return top + body_h

    def chip(self, s, x, y, w, h, txt, color, *, size=10.5, filled=False):
        sh = self.rect(s, x, y, w, h, color if filled else None, line=color, lw=1.0,
                       radius=0.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        self.label(sh, [(txt, size, BG if filled else color, True)])
        return sh

    def chip_row(self, s, x, y, items, color=DIM, *, max_w=CW, size=10.5, h=0.36,
                 gap=0.13, per_char=0.104, pad=0.46):
        """Lay chips left to right, wrapping onto a new row when max_w runs out."""
        cx, cy = x, y
        for txt in items:
            w = per_char * len(txt) + pad
            if cx > x and cx + w > x + max_w:
                cx, cy = x, cy + h + 0.12
            self.chip(s, cx, cy, w, h, txt, color, size=size)
            cx += w + gap
        return cy + h

    def callout(self, s, x, y, w, h, kicker, body, color=AMBER):
        self._check_clear(x, y, w, h, "callout")
        self.rect(s, x, y, w, h, PANEL, line=STROKE, radius=0.06,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        self.rect(s, x, y, 0.055, h, color)
        self.text(s, x + 0.28, y + 0.16, w - 0.5, 0.24, kicker.upper(), 10, color,
                  bold=True)
        self.text(s, x + 0.28, y + 0.46, w - 0.5, h - 0.6, body, 13, TEXT, spacing=1.22)

    def notes(self, s, txt: str):
        s.notes_slide.notes_text_frame.text = txt.strip()

    # -- full-slide templates ----------------------------------------------------------

    def section(self, number: str, title: str, blurb: str, items: list[str]):
        self.module = title
        s = self._blank()
        self.rect(s, 0, 0, 0.28, SH, ACCENT)
        self.rect(s, 0.28, 0, 0.06, SH, CYAN)
        # oversized numeral as a watermark, kept clear of the text column
        self.text(s, 9.1, 1.5, 3.45, 2.7, number, 150, PANEL_2, bold=True, font=UI_LIGHT,
                  align=PP_ALIGN.RIGHT)
        self.text(s, 1.62, 2.72, 7.3, 0.9, title, 44, TEXT, bold=True, font=UI_LIGHT)
        self.text(s, 1.66, 3.74, 7.3, 0.6, blurb, 15, CYAN, spacing=1.2)
        self.chip_row(s, 1.66, 4.52, items, MUTED, max_w=8.8, per_char=0.12, pad=0.44,
                      gap=0.16)
        self._footer(s)
        return s

    def demo(self, s_title: str, cmds: list[str], watch: list[str], fallback: str):
        s = self.slide(s_title, eyebrow="live demo")
        self.chip(s, ML, 1.62, 1.15, 0.34, "▶  DEMO", GREEN, filled=True)
        self.code(s, ML, 2.16, 7.15, cmds, size=13, title="terminal")
        self.card(s, ML + 7.45, 1.62, CW - 7.45, 3.9)
        self.text(s, ML + 7.75, 1.88, CW - 8.05, 0.3, "WHAT TO POINT AT", 10.5, CYAN,
                  bold=True)
        cy = 2.32
        for w in watch:
            self.text(s, ML + 7.75, cy, 0.25, 0.3, "▸", 12, ACCENT)
            self.text(s, ML + 8.05, cy - 0.02, CW - 8.4, 0.7, w, 12.5, TEXT, spacing=1.2)
            cy += 0.62
        self.callout(s, ML, 5.72, CW, 0.95, "if the demo gods say no", fallback, ROSE)
        return s


# --------------------------------------------------------------------------------------
# Deck content
# --------------------------------------------------------------------------------------


WARNINGS: list[str] = []


def build() -> Presentation:
    d = Deck()

    # ---- 01 Title --------------------------------------------------------------------
    s = d._blank()
    d.rect(s, 0, 0, SW, 0.09, ACCENT)
    d.rect(s, 0, 0, 4.4, 0.09, CYAN)
    d.text(s, ML, 1.62, 8.4, 0.34, "MICROSOFT BUILD  ·  DEVELOPER WORKSHOP", 12.5, CYAN,
           bold=True)
    d.text(s, ML, 2.18, 11.0, 1.0, "Microsoft Agent Framework", 54, TEXT, bold=True,
           font=UI_LIGHT)
    d.text(s, ML, 3.24, 11.0, 0.8, "Building Production-Ready AI Agents", 34, ACCENT,
           font=UI_LIGHT)
    d.rect(s, ML, 4.28, 1.6, 0.05, ACCENT)
    d.text(s, ML, 4.62, 9.2, 0.6,
           "A hands-on workshop — every concept below ships as a runnable Python example.",
           15, MUTED, spacing=1.25)
    x = ML
    for tag, col in [("Agents", CYAN), ("Tools", CYAN), ("Memory", CYAN),
                     ("Workflows", VIOLET), ("MCP", VIOLET), ("Middleware", VIOLET),
                     ("Production", GREEN)]:
        w = 0.115 * len(tag) + 0.52
        d.chip(s, x, 5.42, w, 0.38, tag, col)
        x += w + 0.14
    d.rect(s, ML, 6.16, 0.05, 0.62, ACCENT)
    d.text(s, ML + 0.26, 6.14, 6.0, 0.32, "Satendra Kumar", 16, TEXT, bold=True)
    d.text(s, ML + 0.26, 6.48, 8.0, 0.28,
           "agent-framework 1.13  ·  Python 3.13  ·  uv", 11.5, DIM)
    d.notes(s, """
[OPEN]
Welcome. Ninety minutes, eight modules, and by the end you will have written and run a
tool-calling agent, a memory-backed agent, a multi-agent workflow, and a middleware
guardrail — all on your own laptop, against a local model.

[FRAME THE PROMISE]
This is not a "here is the API surface" talk. The framing question all the way through is:
what does it take to move an agent from a notebook demo to something you would put behind
a production endpoint? Every module answers one piece of that.

[PRESENTER TIP]
Before you say a word, have three terminals open and pre-warmed: (1) the vLLM server
already serving, (2) the repo root, (3) a spare. The single most common workshop failure
is a cold model server — first token can take 60+ seconds while weights load.

[LOGISTICS]
Tell people now: clone the repo, run `uv sync`, and start `./vllm-run.sh` in the
background. It downloads weights on first run — get that started during the intro.
""")

    # ---- 02 Goals --------------------------------------------------------------------
    d.module = ""
    s = d.slide("What you'll be able to do after this", eyebrow="outcomes")
    left = [
        ("Explain what the framework is and why it exists",
         "Where it came from, what it replaces, when not to use it."),
        ("Build a tool-calling agent with typed arguments",
         "Function tools, JSON-schema generation, approval gates."),
        ("Give an agent memory that survives the turn",
         "Sessions, context providers, history compaction."),
        ("Orchestrate more than one agent",
         "Functional workflows, graph workflows, concurrent fan-out."),
    ]
    right = [
        ("Extend an agent with external tool servers",
         "MCP over streamable HTTP, transports, trust boundaries."),
        ("Wrap cross-cutting concerns in middleware",
         "Security guardrails, logging, timing, short-circuiting."),
        ("Ship it somewhere real",
         "DevUI, AG-UI, A2A, Azure Functions, OpenTelemetry."),
        ("Recognise the patterns that fail in production",
         "And the ones that hold up under load and cost pressure."),
    ]
    d.bullets(s, ML, 1.78, 5.6, left, size=15, gap=0.98)
    d.bullets(s, ML + 6.0, 1.78, 5.6, right, size=15, gap=0.98)
    d.notes(s, """
[TALK TRACK]
Read these as verbs, not topics. Everyone in this room can already call an LLM API. The
delta we are closing today is everything between "I got a response" and "I would let this
run unattended against a real system."

[KEY POINT]
Call out the last one deliberately. Half the value of a framework is that it names the
patterns — once you can see "this is a context provider" or "this is a middleware
concern", you stop hand-rolling it badly.

[PRESENTER TIP]
Ask for a show of hands: who has shipped an LLM feature to production? Who has shipped an
agent — something that loops and calls tools? The gap between those two numbers is your
whole talk, and it lets you calibrate depth for the rest of the session.
""")

    # ---- 03 Agenda -------------------------------------------------------------------
    s = d.slide("Eight modules", eyebrow="agenda")
    mods = [
        ("01", "Why a framework", "The gap between a completion and an agent", ACCENT),
        ("02", "Core concepts", "Clients, agents, sessions, the run lifecycle", ACCENT),
        ("03", "Tools", "Typed functions, the call loop, approvals", CYAN),
        ("04", "Memory & context", "Sessions, providers, persistence, compaction", CYAN),
        ("05", "MCP", "Standardised tool servers", VIOLET),
        ("06", "Workflows", "Executors, edges, execution, orchestration", VIOLET),
        ("07", "Middleware", "Guardrails and cross-cutting concerns", GREEN),
        ("08", "Production", "Hosting, observability, best practice", GREEN),
    ]
    for i, (num, title, sub, col) in enumerate(mods):
        col_i, row_i = i % 2, i // 2
        x = ML + col_i * 6.1
        y = 1.82 + row_i * 1.22
        c = d.card(s, x, y, 5.7, 1.02)
        d.rect(s, x, y, 0.055, 1.02, col)
        d.text(s, x + 0.3, y + 0.2, 0.7, 0.4, num, 22, col, bold=True, font=UI_LIGHT)
        d.text(s, x + 1.05, y + 0.19, 4.4, 0.32, title, 16, TEXT, bold=True)
        d.text(s, x + 1.05, y + 0.55, 4.4, 0.32, sub, 11.5, MUTED)
    d.notes(s, """
[TALK TRACK]
Four pairs. Modules 1-2 are the mental model. 3-4 are what makes an agent useful — tools
and memory. 5-6 are how you scale beyond one agent. 7-8 are what makes it survivable in
production.

[PACING]
Roughly ten minutes each, but 03 (tools) and 06 (workflows) always run long because those
are where hands go up. Budget for that by keeping 01 tight — the "why" section is the one
you can compress if you are behind.

[PRESENTER TIP]
Say explicitly that modules 5 and 8 are the ones people ask about afterwards, so if
something has to be cut for time you will cut demo repetitions, not those sections. It
buys you permission to move fast without seeming rushed.
""")

    # ---- 04 SECTION 1 ----------------------------------------------------------------
    s = d.section("01", "Why a framework",
                  "What actually breaks when you wire an agent loop by hand",
                  ["the gap", "hand-rolled cost", "lineage", "packages"])
    d.notes(s, """
[TALK TRACK]
Start sceptical on purpose. This audience has seen a lot of agent frameworks, and most of
them are a thin wrapper over a while-loop. If you open with features, you lose them. Open
with the failure modes they have personally hit.

[PRESENTER TIP]
Genuinely ask: "who here has written their own agent loop?" Most hands go up. Then: "who
still maintains it?" Far fewer. That laugh is the argument for this whole section.
""")

    # ---- 05 The gap ------------------------------------------------------------------
    s = d.slide("A completion is not an agent")
    y = 2.05
    lb = d.card(s, ML, y, 5.3, 2.5, fill=PANEL)
    d.text(s, ML + 0.34, y + 0.28, 4.6, 0.3, "CHAT COMPLETION", 11, MUTED, bold=True)
    d.text(s, ML + 0.34, y + 0.68, 4.6, 0.5, "One shot. Stateless.", 20, TEXT, bold=True)
    d.text(s, ML + 0.34, y + 1.22, 4.6, 1.0,
           "prompt in → tokens out. Everything it knows, you put in the prompt. "
           "Everything it does, you do afterwards.", 12.5, MUTED, spacing=1.25)
    d.arrow(s, ML + 5.55, y + 1.1, 0.85, 0.3)
    rb = d.card(s, ML + 6.65, y, 5.15, 2.5, fill=PANEL_2, line=ACCENT)
    d.text(s, ML + 6.99, y + 0.28, 4.4, 0.3, "AGENT", 11, CYAN, bold=True)
    d.text(s, ML + 6.99, y + 0.68, 4.4, 0.5, "Loops. Acts. Remembers.", 20, TEXT,
           bold=True)
    d.text(s, ML + 6.99, y + 1.22, 4.4, 1.0,
           "Decides to call a tool, reads the result, decides again — until the goal is "
           "met or a budget runs out.", 12.5, MUTED, spacing=1.25)
    d.text(s, ML, 5.02, CW, 0.4,
           "The loop is ten lines. Everything around the loop is the hard part.",
           19, TEXT, bold=True, font=UI_LIGHT)
    d.chip_row(s, ML, 5.58,
               ["state across turns", "tool schemas", "retries & timeouts", "approvals",
                "context limits", "traces", "cost caps", "streaming"])
    d.notes(s, """
[TALK TRACK]
Draw the distinction hard: a completion is a pure function; an agent is a process with
state. Once you have a process, you inherit every operational concern that processes have
— lifecycle, observability, failure handling, resource limits.

[KEY POINT]
The chips along the bottom are the real curriculum. Nobody needs a framework to write a
while-loop around a tool call. People need a framework because item eight on that list
arrives three weeks after the demo, and by then the loop has grown to 800 lines.

[PRESENTER TIP]
Pick one chip and tell a war story — "context limits" works well. Describe an agent that
worked beautifully for six turns and then started silently dropping the system prompt.
Concrete beats abstract here.

[TRANSITION]
"So: what would you build if you had to solve all eight of those once, properly?"
""")

    # ---- 06 Hand-rolled cost ---------------------------------------------------------
    s = d.slide("What you end up writing yourself")
    items = [
        ("Message plumbing", "Roles, tool-call IDs, multi-part content, provider drift."),
        ("Schema generation", "Python signatures → JSON Schema, kept in sync by hand."),
        ("The tool loop", "Parse, dispatch, serialise results, guard against loops."),
        ("Conversation state", "Store, load, trim — and decide what 'trim' means."),
        ("Provider swaps", "Each SDK disagrees about streaming, tools, and errors."),
        ("Human approval", "Pausing mid-loop and resuming is genuinely hard."),
        ("Telemetry", "Spans per turn, per tool, per token — retrofitted, always."),
        ("Multi-agent", "Once there are two agents, you have written an orchestrator."),
    ]
    d.bullets(s, ML, 1.82, 5.5, items[:4], gap=0.92)
    d.bullets(s, ML + 6.0, 1.82, 5.5, items[4:], gap=0.92)
    d.callout(s, ML, 5.72, CW, 0.98, "the actual argument",
              "None of these are hard problems. They are all boring problems — and you "
              "will solve each one slightly differently in every service you own.")
    d.notes(s, """
[TALK TRACK]
Walk this list fast — it should feel like a slightly painful inventory rather than a
lecture. The audience is checking off the ones they have personally written.

[KEY POINT]
Land the callout: the case for a framework here is not capability, it is consistency.
Eight teams hand-rolling eight tool loops is eight different retry semantics and eight
different bugs in production.

[HONEST CAVEAT]
Say out loud that if all you need is one prompt and one function call, you do not need
this. Frameworks earn their keep at the second agent, the second provider, or the first
compliance review. Being honest here buys credibility for the rest of the deck.

[PRESENTER TIP]
"Provider swaps" is the one architects react to. If the room skews architect, dwell there
and mention that the same Agent code in this repo runs against OpenAI, Azure, Foundry,
Bedrock, Gemini, Ollama and a local vLLM server.
""")

    # ---- 07 What is MAF --------------------------------------------------------------
    s = d.slide("Microsoft Agent Framework")
    d.text(s, ML, 1.78, 11.4, 0.6,
           "One SDK and runtime for agents and multi-agent workflows — Python and .NET.",
           19, TEXT, spacing=1.2, font=UI_LIGHT)
    y = 2.62
    a = d.card(s, ML, y, 3.55, 1.55, fill=PANEL)
    d.label(a, [("AutoGen", 19, TEXT, True),
                ("research velocity", 12, MUTED, False),
                ("simple multi-agent abstractions", 11, DIM, False)])
    b = d.card(s, ML, y + 1.78, 3.55, 1.55, fill=PANEL)
    d.label(b, [("Semantic Kernel", 19, TEXT, True),
                ("enterprise hardening", 12, MUTED, False),
                ("state, telemetry, connectors", 11, DIM, False)])
    d.arrow(s, ML + 3.78, y + 1.42, 0.9, 0.32)
    c = d.card(s, ML + 4.95, y, 3.5, 3.33, fill=PANEL_2, line=ACCENT)
    d.label(c, [("Agent Framework", 22, CYAN, True), (" ", 8, MUTED, False),
                ("one object model", 13, TEXT, False),
                ("one runtime", 13, TEXT, False),
                ("graph workflows", 13, TEXT, False),
                ("open source", 13, TEXT, False)])
    adds = [("Durable workflows", "typed graph, checkpoints"),
            ("Open standards", "MCP · A2A · AG-UI · OTel"),
            ("Provider-neutral", "swap the client, keep the agent")]
    cy = y
    for head, sub in adds:
        card = d.card(s, ML + 8.72, cy, 3.1, 1.02)
        d.label(card, [(head, 13.5, TEXT, True), (sub, 10.5, MUTED, False)])
        cy += 1.16
    d.notes(s, """
[TALK TRACK]
The lineage matters because it explains the design. AutoGen brought the research-side
ergonomics — spinning up multi-agent conversations in a few lines. Semantic Kernel brought
the parts enterprises actually asked for — state management, telemetry, connectors, a
plugin model. Agent Framework is the convergence, plus a real workflow engine.

[KEY POINT]
If someone asks "is AutoGen dead / should I migrate off SK" — the honest framing is that
this is the forward-looking line, both prior projects informed it, and migration guidance
exists in the docs. Do not oversell a hard deprecation you cannot substantiate.

[PRESENTER TIP]
The right-hand column is what is genuinely new rather than merged. Emphasise open
standards: MCP for tools, A2A for agent-to-agent, AG-UI for front ends, OpenTelemetry for
traces. That is an unusually standards-forward posture and it is a real differentiator.
""")

    # ---- 08 Packages -----------------------------------------------------------------
    s = d.slide("How the packages are laid out")
    core = d.card(s, ML, 1.85, CW, 1.18, fill=PANEL_2, line=ACCENT)
    d.text(s, ML + 0.34, 2.05, 4.0, 0.32, "agent_framework", 19, CYAN, bold=True,
           font=MONO)
    d.text(s, ML + 0.34, 2.44, 11.2, 0.4,
           "Agent · tool · AgentSession · ContextProvider · WorkflowBuilder · Executor · "
           "middleware · compaction", 12.5, MUTED)
    groups = [
        ("Model providers", CYAN,
         ["openai", "azure", "foundry", "anthropic", "gemini", "bedrock", "mistral",
          "ollama", "claude", "copilotstudio"]),
        ("Memory & stores", VIOLET, ["redis", "mem0", "azure_cosmos", "azure_ai_search"]),
        ("Surfaces & hosting", GREEN,
         ["devui", "ag_ui", "a2a", "azurefunctions", "durabletask", "chatkit",
          "foundry_hosting"]),
    ]
    x = ML
    for name, col, pkgs in groups:
        w = 3.82
        d.card(s, x, 3.28, w, 2.35)
        d.rect(s, x, 3.28, w, 0.055, col)
        d.text(s, x + 0.28, 3.5, w - 0.5, 0.3, name, 14, TEXT, bold=True)
        cy = 3.9
        for i in range(0, len(pkgs), 2):
            row = "  ·  ".join(pkgs[i : i + 2])
            d.text(s, x + 0.28, cy, w - 0.5, 0.26, row, 11.5, MUTED, font=MONO)
            cy += 0.3
        x += w + 0.2
    d.callout(s, ML, 5.86, CW, 0.88, "install note",
              "The agent-framework meta-package pulls every provider in — this repo already "
              "has all of the above installed, so new examples need no dependency edits.",
              CYAN)
    d.notes(s, """
[TALK TRACK]
Core is small and provider-agnostic on purpose. Everything vendor-specific lives in a
sibling distribution. That is why `Agent` never mentions OpenAI — it takes a client.

[GOTCHA WORTH SAYING OUT LOUD]
The same class is reachable two ways: `from agent_framework_openai import ...` and
`from agent_framework.openai import ...`. Both work — namespace package plus per-provider
distribution. The examples in this repo mix both, and that is not a bug. If you see it in
review, do not "fix" it blindly.

[PRESENTER TIP]
If anyone asks about install size — yes, the meta-package is heavy. In production you pin
the specific provider distributions you need. For a workshop, the meta-package means
nobody gets blocked on a missing import halfway through.
""")

    # ---- 09 SECTION 2 ----------------------------------------------------------------
    s = d.section("02", "Core concepts",
                  "Five primitives, and the lifecycle that connects them",
                  ["chat client", "agent", "session", "context provider", "lifecycle"])
    d.notes(s, """
[TALK TRACK]
This is the module that pays for the rest of the day. If people leave with the run
lifecycle diagram in their heads, everything afterwards — tools, memory, middleware — is
just "which hook does this hang off".

[PRESENTER TIP]
Slow down here. It is tempting to rush to the tool demo because that is the fun part, but
a shaky mental model at module 2 means confused questions at module 7.
""")

    # ---- 10 Five primitives ----------------------------------------------------------
    s = d.slide("Five primitives")
    prims = [
        ("Chat client", "OpenAIChatCompletionClient", "Talks to a model. Swappable.",
         CYAN),
        ("Agent", "Agent(...)", "Instructions + tools + client. The unit of behaviour.",
         ACCENT),
        ("Tool", "@tool", "A typed Python function the model may call.", VIOLET),
        ("Session", "agent.create_session()", "Carries conversation state across runs.",
         AMBER),
        ("Context provider", "ContextProvider", "Injects and captures context per run.",
         GREEN),
    ]
    w = 2.24
    for i, (name, api, desc, col) in enumerate(prims):
        x = ML + i * (w + 0.16)
        d.card(s, x, 2.0, w, 3.05)
        d.rect(s, x, 2.0, w, 0.055, col)
        d.rect(s, x + w / 2 - 0.22, 2.32, 0.44, 0.44, None, line=col, lw=1.5,
               shape=MSO_SHAPE.OVAL)
        d.text(s, x, 2.42, w, 0.3, str(i + 1), 15, col, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, x + 0.2, 2.96, w - 0.4, 0.34, name, 15, TEXT, bold=True,
               align=PP_ALIGN.CENTER)
        d.text(s, x + 0.12, 3.36, w - 0.24, 0.3, api, 9.5, col, font=MONO,
               align=PP_ALIGN.CENTER)
        d.text(s, x + 0.2, 3.78, w - 0.4, 1.1, desc, 11.5, MUTED, spacing=1.2,
               align=PP_ALIGN.CENTER)
    d.callout(s, ML, 5.42, CW, 1.0, "composition, not inheritance",
              "You never subclass Agent. You compose it: hand it a different client, "
              "another tool, one more context provider, a middleware. That is the whole "
              "extensibility story.", CYAN)
    d.notes(s, """
[TALK TRACK]
Five nouns. Everything in the framework is one of these or a composition of them.

[KEY POINT]
The callout is the design principle worth repeating. There is no `class MyAgent(Agent)` in
idiomatic code. When someone asks "how do I customise X", the answer is almost always
"pass something in", not "override something".

[ORDER MATTERS]
Introduce them in dependency order: a client is useless alone, an agent needs a client, a
session needs an agent, a context provider hangs off the agent and reads/writes session
state. Tools are the only one that stands alone as plain Python.

[PRESENTER TIP]
Ask the room to predict which of the five is the one people misuse most. Answer: sessions
— specifically, sharing one session across concurrent users. Flag it now, revisit in
module 8.
""")

    # ---- 11 Chat clients -------------------------------------------------------------
    s = d.slide("Chat clients: the model is a parameter")
    d.code(s, ML, 1.8, 6.9, [
        "from agent_framework import Agent",
        "from agent_framework_openai import OpenAIChatCompletionClient",
        "",
        "agent = Agent(",
        "    client=OpenAIChatCompletionClient(),   # swap this line only",
        "    name=\"HelloAgent\",",
        "    instructions=\"You are a friendly assistant. Be brief.\",",
        ")",
    ], title="the only provider-specific line")
    d.card(s, ML + 7.2, 1.8, 4.6, 2.2)
    d.text(s, ML + 7.5, 2.02, 4.0, 0.3, "SAME AGENT, OTHER BACKENDS", 10.5, CYAN,
           bold=True)
    cy = 2.42
    for a, b in [("AzureOpenAIChatClient", "Azure OpenAI"),
                 ("FoundryChatClient", "Azure AI Foundry"),
                 ("AnthropicChatClient", "Claude"),
                 ("OllamaChatClient", "local Ollama")]:
        d.text(s, ML + 7.5, cy, 2.7, 0.26, a, 10.5, TEXT, font=MONO)
        d.text(s, ML + 10.3, cy, 1.4, 0.26, b, 10.5, MUTED)
        cy += 0.36
    d.callout(s, ML, 4.86, CW, 1.18, "the one distinction that bites people",
              "OpenAIChatCompletionClient → Chat Completions API.   "
              "OpenAIChatClient → the newer Responses API.\n"
              "Local servers such as vLLM only implement Chat Completions. Pick the wrong "
              "one and you get a 404 that looks like an auth failure.", ROSE)
    d.notes(s, """
[TALK TRACK]
The point of this slide is one line of code. `client=` is the only place the provider
appears. Everything downstream — tools, sessions, workflows, middleware — is written
against the abstraction.

[THE GOTCHA — say this twice]
`OpenAIChatCompletionClient` and `OpenAIChatClient` are different clients for different
OpenAI APIs. The workshop runs against a local vLLM server, which serves Chat Completions
only. Every example here uses the ChatCompletion variant deliberately.

[PRESENTER TIP]
This is the single highest-value troubleshooting fact in the deck. When someone's laptop
throws a 404 in twenty minutes, it is this. Consider writing it on a whiteboard and
leaving it up.

[ARCHITECT ANGLE]
Provider-neutrality is not just portability theatre — it is how you run cheap local models
in CI and a frontier model in production without branching your agent code.
""")

    # ---- 12 Hello agent code ---------------------------------------------------------
    s = d.slide("Your first agent", eyebrow="examples/01_hello_agent.py")
    d.code(s, ML, 1.8, 7.35, [
        "import asyncio",
        "from agent_framework import Agent",
        "from agent_framework_openai import OpenAIChatCompletionClient",
        "from dotenv import load_dotenv",
        "",
        "load_dotenv()",
        "",
        "async def main() -> None:",
        "    agent = Agent(",
        "        client=OpenAIChatCompletionClient(),",
        "        name=\"HelloAgent\",",
        "        instructions=\"Keep your answers brief.\",",
        "    )",
        "    result = await agent.run(\"What is the capital of France?\")",
        "    print(result)",
        "",
        "asyncio.run(main())",
    ], size=13)
    notes_items = [
        ("Async everywhere", "agent.run() is a coroutine. No sync shim."),
        ("Config from env", "load_dotenv() feeds base URL, key and model."),
        ("Instructions ≠ prompt", "System-level, applied to every run."),
        ("Result is rich", "result.text, .messages, .value — not just a string."),
    ]
    d.bullets(s, ML + 7.75, 1.9, 4.05, notes_items, size=14, gap=0.92)
    d.notes(s, """
[TALK TRACK]
Seventeen lines, and four of them are imports. Point out that there is no client
configuration in sight — base URL, key and model all come from the environment via
`load_dotenv()`.

[WALK THE CODE]
Agent takes a client, a name, and instructions. `name` matters more than it looks: it
shows up in traces, in multi-agent transcripts as `author_name`, and in DevUI. Give agents
real names, not "agent1".

[KEY POINT]
`instructions` is not the user prompt. It is applied to every run in the session. A common
beginner bug is stuffing per-request context into instructions instead of into the message
or a context provider.

[DEMO CUE]
uv run python examples/01_hello_agent.py

[PRESENTER TIP]
Run it twice — once as written, once after changing instructions to something absurd like
"answer only in haiku". Nothing lands the instructions/prompt distinction faster.
""")

    # ---- 13 Lifecycle ----------------------------------------------------------------
    s = d.slide("The run lifecycle")
    stages = [("agent.run()", ACCENT), ("context providers\nbefore_run", GREEN),
              ("agent\nmiddleware", VIOLET), ("model call", CYAN),
              ("context providers\nafter_run", GREEN)]
    x, w = ML, 1.98
    for i, (t, col) in enumerate(stages):
        sh = d.card(s, x, 2.12, w, 1.05, fill=PANEL, line=col)
        d.label(sh, [(l, 12, TEXT, True) for l in t.split("\n")])
        if i < len(stages) - 1:
            d.arrow(s, x + w + 0.06, 2.52, 0.35, 0.26, col)
        x += w + 0.47
    # the tool loop hangs off the model call and cycles back into it
    model_cx = ML + 3 * (w + 0.47) + w / 2
    d.arrow(s, model_cx - 0.42, 3.28, 0.24, 0.5, AMBER, shape=MSO_SHAPE.DOWN_ARROW)
    d.arrow(s, model_cx + 0.18, 3.28, 0.24, 0.5, AMBER, shape=MSO_SHAPE.UP_ARROW)
    loop = d.card(s, model_cx - 2.15, 3.86, 4.3, 1.0, fill=PANEL, line=AMBER)
    d.label(loop, [("tool call requested?", 12.5, AMBER, True),
                   ("function middleware → tool → result → model again", 10.5, MUTED,
                    False)])
    d.text(s, model_cx - 2.15, 4.96, 4.3, 0.3,
           "repeats until the model stops asking for tools", 11.5, MUTED,
           align=PP_ALIGN.CENTER)
    d.callout(s, ML, 5.18, CW, 1.05, "read the diagram as an onion",
              "before_run injects context → middleware wraps the call → the model may loop "
              "through tools → after_run captures what happened. Every extension point you "
              "will meet today is one of those four rings.", CYAN)
    d.notes(s, """
[TALK TRACK]
This is the slide to draw on a whiteboard rather than read. Four rings around one model
call.

[WALK IT]
1. You call agent.run(). 2. Every registered context provider gets before_run — this is
where memory injects "the user's name is Alice" into the instructions. 3. Agent middleware
wraps the whole invocation and can short-circuit it entirely. 4. The model responds — and
if it asks for a tool, function middleware wraps that call, the tool runs, the result goes
back, and the model gets another turn. That inner loop repeats. 5. after_run fires so
providers can persist anything they learned.

[KEY POINT]
The tool loop is inside the middleware, not outside it. That is why one agent middleware
sees the whole multi-turn tool exchange as a single unit, while function middleware fires
per tool call. People get this backwards constantly.

[PRESENTER TIP]
Ask: "where would you put a token budget check?" Good answers differ — agent middleware
for a per-request cap, function middleware for a per-tool cap. Let the room argue for
thirty seconds; it proves they have the model.
""")

    # ---- 14 Streaming ----------------------------------------------------------------
    s = d.slide("Streaming, options and structured results")
    d.code(s, ML, 1.8, 7.35, [
        "# non-streaming — one await, full response",
        "result = await agent.run(\"What is the capital of France?\")",
        "",
        "# per-call options ride alongside the prompt",
        "result = await agent.run(prompt, options={",
        "    \"temperature\": 0.3,",
        "    \"max_tokens\": 150,",
        "})",
        "",
        "# streaming — same method, stream=True",
        "async for chunk in agent.run(prompt, stream=True):",
        "    if chunk.text:",
        "        print(chunk.text, end=\"\", flush=True)",
    ], size=13)
    d.bullets(s, ML + 7.75, 1.9, 4.05, [
        ("One method, two shapes", "stream=True returns an async iterator."),
        ("Guard chunk.text", "Chunks may carry tool calls, not text."),
        ("options is per call", "default_options on Agent sets the baseline."),
        ("Cap tokens deliberately", "The workshop model has an 8k window."),
    ], size=14, gap=0.92)
    d.notes(s, """
[TALK TRACK]
Same `run` method either way — `stream=True` flips it from awaitable to async iterator.
That symmetry is deliberate and it means you can develop non-streaming and ship streaming
without restructuring.

[WATCH FOR]
`if chunk.text:` is not defensive noise. Streaming updates carry more than text — tool call
deltas, usage, annotations. Iterating naively and printing `chunk` prints objects.

[REPO SPECIFIC]
Examples here pass `max_tokens` explicitly because the local vLLM server runs with
`--max-model-len=8192`. On a hosted frontier model you would usually leave it off. Mention
that `default_options={"max_tokens": 500}` on the Agent sets a baseline for every call —
see examples/08.

[PRESENTER TIP]
If the room is UI-oriented, note that streaming plus AG-UI (module 8) is how you get a
token-by-token chat front end with no bespoke transport code.
""")

    # ---- 15 DEMO 1 -------------------------------------------------------------------
    s = d.demo(
        "Demo — hello agent, streaming, options",
        ["# start the model server first (separate terminal)",
         "./vllm-run.sh",
         "",
         "# then:",
         "uv run python examples/01_hello_agent.py",
         "",
         "# no model needed — good smoke test:",
         "uv run python examples/11_graph_workflow.py"],
        ["The first response is blocking — the whole answer lands at once.",
         "The second passes temperature and max_tokens per call.",
         "The third streams: tokens appear as they are produced.",
         "Edit instructions live and re-run — show the behaviour change."],
        "Run examples/11_graph_workflow.py — it needs no model at all and still "
        "demonstrates the framework executing. Keep talking while vLLM warms up.")
    d.notes(s, """
[SETUP]
vLLM must already be serving on localhost:8000 before this slide. First run downloads
Qwen3.5-4B weights — do not do that live.

[DEMO SCRIPT]
1. Show the file (17 lines). 2. Run it. 3. Point at the three outputs and map each to the
code. 4. Change instructions to "answer only in haiku", re-run, get a laugh, move on.

[PRESENTER TIP]
Keep your terminal font large — 18pt minimum. And pre-clear scrollback; a wall of vLLM
debug logging above your output kills the moment.

[COMMON FAILURE]
404 from the client almost always means OPENAI_BASE_URL is unset or you are on
OpenAIChatClient instead of OpenAIChatCompletionClient. Have .env open in a tab.

[TIMING]
Three minutes. Do not let this one sprawl — the interesting demos are later.
""")

    # ---- 16 SECTION 3 ----------------------------------------------------------------
    s = d.section("03", "Tools",
                  "Giving the model a typed, governed way to touch your systems",
                  ["@tool", "schema", "the loop", "approvals"])
    d.notes(s, """
[TALK TRACK]
Tools are where an agent stops being a chatbot. Everything before this was text in, text
out. From here the model can cause effects — which is exactly why the approval slide at
the end of this module matters more than the syntax at the start.

[PRESENTER TIP]
This module always overruns because it is the most immediately useful. Watch the clock.
""")

    # ---- 17 Tool anatomy -------------------------------------------------------------
    s = d.slide("Anatomy of a tool")
    d.code(s, ML, 1.8, 7.1, [
        "@tool(approval_mode=\"never_require\")",
        "def get_weather(",
        "    location: Annotated[str, Field(description=\"City to look up.\")],",
        ") -> str:",
        "    \"\"\"Get the weather for a given location.\"\"\"",
        "    return f\"{location}: sunny, high of 24C\"",
    ], size=12)
    d.text(s, ML, 3.82, 7.1, 0.3, "WHAT THE MODEL ACTUALLY RECEIVES", 10.5, CYAN,
           bold=True)
    maps = [("function name", "get_weather", "the callable identifier"),
            ("docstring", "\"Get the weather…\"", "when to call it"),
            ("type hints", "location: str", "types → JSON Schema"),
            ("Field(description=…)", "\"City to look up.\"", "what goes in the argument")]
    cy = 4.18
    for a, b, c in maps:
        d.text(s, ML, cy, 2.2, 0.28, a, 12, TEXT, bold=True)
        d.arrow(s, ML + 2.28, cy + 0.05, 0.28, 0.18, DIM)
        d.text(s, ML + 2.7, cy, 2.3, 0.28, b, 11.5, AMBER, font=MONO)
        d.text(s, ML + 5.1, cy, 2.2, 0.28, c, 11, MUTED)
        cy += 0.40
    d.card(s, ML + 7.4, 1.8, 4.4, 3.0, fill=PANEL_2, line=STROKE)
    d.text(s, ML + 7.7, 2.02, 3.8, 0.3, "THE DOCSTRING IS AN API", 10.5, AMBER, bold=True)
    d.text(s, ML + 7.7, 2.42, 3.8, 2.2,
           "It is not a comment for humans. It is the only thing telling the model when "
           "this tool is the right choice.\n\nVague docstring → the model guesses. Guessing "
           "looks like a flaky agent, and you will debug the prompt for a day before you "
           "look here.", 12.5, TEXT, spacing=1.25)
    d.callout(s, ML, 5.76, CW, 0.98, "design rule",
              "Write tools the way you would write a public API for a junior engineer with "
              "no context: narrow scope, unambiguous name, explicit units, errors returned "
              "as readable strings the model can act on.", GREEN)
    d.notes(s, """
[TALK TRACK]
The decorator does one job: it turns a Python signature into a JSON Schema the model can
be trained on. Name, docstring, type hints and Field descriptions are all inputs to that
schema.

[KEY POINT — the docstring]
Say this plainly: your docstring is prompt engineering. It is the highest-leverage text in
the whole agent, and it is the one most teams write last and worst.

[PRACTICAL ADVICE]
Return errors as strings, not exceptions, when you want the model to recover — "no city
found named Xyz, try a fuller name" gives the model something to do. Raise when you want
the run to fail loudly.

[UNITS]
Tell the war story about units. A tool that returns "24" versus "24C" versus "the high is
24 degrees Celsius" produces measurably different downstream behaviour.

[PRESENTER TIP]
Ask the room to critique a deliberately bad docstring — "gets data" — before you show the
good one. Participation here wakes people up mid-session.
""")

    # ---- 18 Tool loop ----------------------------------------------------------------
    s = d.slide("The tool-call loop")
    flow = [("user message", ACCENT), ("agent", CYAN), ("model", VIOLET),
            ("tool", AMBER)]
    fw, fgap = 2.3, 0.55
    x = ML
    for i, (t, col) in enumerate(flow):
        sh = d.card(s, x, 2.05, fw, 0.9, fill=PANEL, line=col)
        d.label(sh, [(t, 13, TEXT, True)])
        if i < len(flow) - 1:
            d.arrow(s, x + fw + 0.08, 2.39, 0.39, 0.24, col)
        x += fw + fgap
    # the result travels back to the model, which then decides again
    d.arrow(s, ML + 6.3, 3.08, 3.2, 0.22, AMBER, shape=MSO_SHAPE.LEFT_ARROW)
    d.text(s, ML + 6.3, 3.34, 3.2, 0.26, "result", 11, AMBER, align=PP_ALIGN.CENTER)
    steps = [
        ("1", "Model returns a tool call", "name + JSON arguments, not prose", CYAN),
        ("2", "Framework validates & dispatches", "schema check, then your Python runs",
         VIOLET),
        ("3", "Result goes back as a message", "the model now has new facts", AMBER),
        ("4", "Model decides again", "answer, or call another tool — repeat", GREEN),
    ]
    for i, (n, head, sub, col) in enumerate(steps):
        y = 3.85 + i * 0.62
        d.rect(s, ML, y, 0.34, 0.34, None, line=col, lw=1.3, shape=MSO_SHAPE.OVAL)
        d.text(s, ML, y + 0.05, 0.34, 0.26, n, 12, col, bold=True, align=PP_ALIGN.CENTER)
        d.text(s, ML + 0.56, y + 0.01, 4.4, 0.3, head, 14, TEXT, bold=True)
        d.text(s, ML + 5.1, y + 0.04, 3.6, 0.3, sub, 12, MUTED)
    d.card(s, ML + 8.9, 3.85, 2.9, 2.35, fill=PANEL_2, line=ROSE)
    d.text(s, ML + 9.18, 4.06, 2.4, 0.3, "BOUND THE LOOP", 10.5, ROSE, bold=True)
    d.text(s, ML + 9.18, 4.44, 2.4, 1.7,
           "Nothing stops a model calling a tool forever except limits you set: max "
           "iterations, token budget, timeouts, and idempotent tools.", 12, TEXT,
           spacing=1.25)
    d.notes(s, """
[TALK TRACK]
Emphasise that step 1 is structured output, not text parsing. The model emits a tool call
with JSON arguments; the framework validates against the schema before your function ever
sees it. That validation boundary is one of the quiet wins over hand-rolled loops.

[KEY POINT]
Step 4 is the recursion. The loop terminates when the model stops asking for tools — which
means termination is a model decision unless you impose limits. Hence the red panel.

[WAR STORY]
Describe a retrieval tool that returned "no results found", the model rephrased and called
again, and again — forty times, until the context window filled. Fix was a call budget
plus a tool that said "no results; do not retry with a synonym".

[PRESENTER TIP]
If you have time, ask what happens when two tool calls come back in one response. Answer:
they can be dispatched in parallel — which is why your tools should be safe to run
concurrently and ideally idempotent.
""")

    # ---- 19 Tools code ---------------------------------------------------------------
    s = d.slide("Tools in practice", eyebrow="examples/02_add_tools.py")
    d.code(s, ML, 1.8, 7.5, [
        "from agent_framework import Agent, tool",
        "from pydantic import Field",
        "",
        "@tool(approval_mode=\"never_require\")",
        "def get_weather(",
        "    location: Annotated[str, Field(description=\"The location.\")],",
        ") -> str:",
        "    \"\"\"Get the weather for a given location.\"\"\"",
        "    conditions = [\"sunny\", \"cloudy\", \"rainy\", \"stormy\"]",
        "    return f\"{location} is {choice(conditions)}, high {randint(10, 30)}C\"",
        "",
        "agent = Agent(",
        "    client=OpenAIChatCompletionClient(),",
        "    name=\"WeatherAgent\",",
        "    instructions=\"Use get_weather to answer.\",",
        "    tools=[get_weather],",
        ")",
    ], size=12.5)
    d.bullets(s, ML + 7.9, 1.9, 3.9, [
        ("tools= takes a list", "Or a single tool — both accepted."),
        ("Plain functions work too", "The decorator adds approval control."),
        ("Instructions steer usage", "Name the tool if the model hesitates."),
        ("Async tools supported", "Use async def for I/O-bound work."),
    ], size=13.5, gap=0.86)
    d.notes(s, """
[TALK TRACK]
This is example 02 essentially verbatim. Note that `tools=` accepts a list or a bare tool —
example 16 passes `tools=get_weather` with no list, and that is fine.

[SUBTLE POINT]
You do not strictly need the `@tool` decorator — a plain typed function with a docstring
works (see dev_ui/agent.py). You reach for the decorator when you want to control approval
behaviour or other tool-level settings.

[MODEL BEHAVIOUR]
With a small local model, tool selection is less reliable than with a frontier model. That
is why the instructions here explicitly say "Use get_weather to answer". Mention that this
is a real production technique for cheap models, not a workshop hack.

[DEMO CUE]
uv run python examples/02_add_tools.py

[PRESENTER TIP]
Ask it about a city, then ask something unrelated — "what is 2+2" — and show that the model
does not call the tool. Demonstrating restraint is more convincing than demonstrating use.
""")

    # ---- 20 Approvals ----------------------------------------------------------------
    s = d.slide("Approvals: the difference between a demo and production")
    a = d.card(s, ML, 1.9, 5.6, 2.35, fill=PANEL, line=STROKE)
    d.rect(s, ML, 1.9, 5.6, 0.055, GREEN)
    d.text(s, ML + 0.34, 2.14, 4.9, 0.3, "never_require", 17, TEXT, bold=True, font=MONO)
    d.text(s, ML + 0.34, 2.56, 4.9, 1.5,
           "Runs immediately. Correct for reads: lookups, search, computation, anything "
           "you would happily let a stranger call a thousand times.", 13, MUTED,
           spacing=1.25)
    d.chip(s, ML + 0.34, 3.62, 1.5, 0.34, "safe reads", GREEN)
    b = d.card(s, ML + 6.2, 1.9, 5.6, 2.35, fill=PANEL_2, line=AMBER)
    d.rect(s, ML + 6.2, 1.9, 5.6, 0.055, AMBER)
    d.text(s, ML + 6.54, 2.14, 4.9, 0.3, "always_require", 17, TEXT, bold=True, font=MONO)
    d.text(s, ML + 6.54, 2.56, 4.9, 1.5,
           "Pauses the run and surfaces an approval request. Correct for writes, "
           "payments, emails, deletes — anything with a blast radius.", 13, TEXT,
           spacing=1.25)
    d.chip(s, ML + 6.54, 3.62, 1.9, 0.34, "effects & spend", AMBER)
    d.text(s, ML, 4.52, CW, 0.34,
           "Every example in this repo uses never_require for brevity — and says so in a "
           "comment.", 14, ROSE, bold=True)
    d.callout(s, ML, 5.06, CW, 1.22, "how to decide",
              "Ask one question: if the model called this tool with the worst plausible "
              "arguments, ten times in a row, at 3am, unattended — what happens?\n"
              "Nothing → never_require.   Anything → always_require, and make the tool "
              "idempotent as well.", ROSE)
    d.notes(s, """
[TALK TRACK]
This is the slide that makes architects relax. Human-in-the-loop is not bolted on; pausing
and resuming mid-loop is a first-class capability, which is genuinely hard to build
yourself because you have to serialise the pending tool call and the conversation state.

[BE HONEST ABOUT THE REPO]
Every sample here sets never_require, and the framework's own samples carry the same
comment. That is a teaching shortcut. Say clearly: do not copy that line into a service
that spends money.

[THE 3AM TEST]
Give the room the heuristic in the callout and then apply it live to three tools they
suggest. It works better as an interaction than as a bullet.

[FOLLOW-UP]
Approval pairs with sessions — the approval round-trip needs somewhere to park state.
Point forward to module 4.
""")

    # ---- 21 DEMO 2 -------------------------------------------------------------------
    s = d.demo(
        "Demo — tools and the call loop",
        ["uv run python examples/02_add_tools.py",
         "",
         "# then try, live:",
         "#  - ask about two cities in one message",
         "#  - ask something with no tool at all",
         "#  - break the docstring and watch selection degrade"],
        ["The model emits a tool call, not prose — show the arguments.",
         "Two cities in one prompt can produce two tool calls.",
         "Off-topic questions should NOT trigger the tool.",
         "Vandalise the docstring to \"gets data\" and re-run."],
        "The weather tool is pure Python with randomness — it always returns something. "
        "If the model refuses to call it, restate the instruction to name the tool "
        "explicitly and re-run.")
    d.notes(s, """
[DEMO SCRIPT]
The docstring vandalism is the memorable one. Change the docstring to "gets data", re-run,
and the small local model will often stop calling the tool or call it with a nonsense
argument. Restore it and it works again. That is the whole "docstring is an API" lesson in
forty seconds.

[PRESENTER TIP]
Have the edit ready in a second editor tab so you are not typing under pressure. Use git
to restore: `git checkout examples/02_add_tools.py`.

[IF ASKED ABOUT PARALLEL CALLS]
Two-city prompts sometimes yield two tool calls in one response, sometimes sequential ones
— depends on the model. Do not promise a specific behaviour live.

[TIMING]
Four minutes.
""")

    # ---- 22 SECTION 4 ----------------------------------------------------------------
    s = d.section("04", "Memory & context",
                  "Three horizons: the turn, the session, and everything before",
                  ["sessions", "providers", "persistence", "compaction"])
    d.notes(s, """
[TALK TRACK]
Memory is where most agent projects quietly go wrong. Not because it is conceptually hard,
but because "just append to the message list" works perfectly until turn thirty, and then
it fails as a cost problem and a quality problem simultaneously.

[PRESENTER TIP]
Ask how many people have hit a context-length error in production. Then ask how many
solved it by truncating the oldest messages. Then ask how that went.
""")

    # ---- 23 Three horizons -----------------------------------------------------------
    s = d.slide("Three horizons of memory")
    tiers = [
        ("The turn", "message list", ACCENT,
         "What you pass into a single run. Always present, never persisted."),
        ("The session", "AgentSession", CYAN,
         "Conversation state across runs. In-memory by default; pluggable stores."),
        ("Beyond the session", "ContextProvider + store", VIOLET,
         "Facts, preferences, prior conversations — retrieved and injected per run."),
    ]
    for i, (name, api, col, desc) in enumerate(tiers):
        x = ML + i * 4.0
        d.card(s, x, 1.95, 3.78, 2.5)
        d.rect(s, x, 1.95, 3.78, 0.055, col)
        d.text(s, x + 0.3, 2.2, 3.2, 0.34, name, 18, TEXT, bold=True)
        d.text(s, x + 0.3, 2.62, 3.2, 0.28, api, 11.5, col, font=MONO)
        d.text(s, x + 0.3, 3.06, 3.2, 1.2, desc, 12.5, MUTED, spacing=1.25)
        d.rect(s, x + 0.3, 4.06, 0.5 + i * 1.1, 0.16, col)
    d.text(s, ML, 4.62, CW, 0.3, "grows →", 11, DIM)
    d.code(s, ML, 4.98, CW, [
        "session = agent.create_session()          # one conversation",
        "await agent.run(\"My name is Alice.\", session=session)",
        "await agent.run(\"What do you remember?\", session=session)   # knows",
    ], size=13)
    d.notes(s, """
[TALK TRACK]
Three horizons, three different mechanisms, and the mistake is using one for another's job.

[THE TURN]
The message list. Cheap, explicit, gone afterwards. Anything you can recompute belongs
here, not in memory.

[THE SESSION]
`agent.create_session()` returns a handle you pass to every run. It carries history plus a
state dict that context providers read and write. In-memory by default — which is fine for
a CLI and wrong for a web service.

[BEYOND]
This is where retrieval lives. A context provider that queries Redis, mem0, Cosmos or
Azure AI Search and injects the top few facts into the instructions for this run.

[KEY POINT]
Sessions are not a database. If your process restarts, an in-memory session is gone. For
anything user-facing, back it with a real history provider — module 8 revisits this.

[PRESENTER TIP]
The three-bar visual maps to cost. Longer horizon, more retrieval, more tokens. Every
memory decision is a cost decision.
""")

    # ---- 24 Sessions code ------------------------------------------------------------
    s = d.slide("Sessions: multi-turn in three lines", eyebrow="examples/03_multi_turn.py")
    d.code(s, ML, 1.85, 7.4, [
        "agent = Agent(",
        "    client=OpenAIChatCompletionClient(),",
        "    name=\"ConversationAgent\",",
        "    instructions=\"Keep your answers brief.\",",
        ")",
        "",
        "session = agent.create_session()",
        "",
        "await agent.run(\"My name is Alice and I love hiking.\", session=session)",
        "",
        "# same session → the agent still has the earlier turns",
        "result = await agent.run(\"What do you remember about me?\",",
        "                         session=session)",
    ], size=12.5)
    d.bullets(s, ML + 7.8, 1.95, 4.0, [
        ("Omit session → amnesia", "Each run starts clean. Sometimes that's what you want."),
        ("One session per conversation", "Never share one across users or requests."),
        ("session.state is a dict", "Providers namespace their data by source_id."),
        ("Swap the backing store", "InMemory, Redis, Cosmos — same API."),
    ], size=13.5, gap=0.94)
    d.notes(s, """
[TALK TRACK]
Deliberately unremarkable code — that is the point. The session is a parameter, not a
global, not an ambient context.

[KEY POINT — the production bug]
"One session per conversation." The failure mode is a service that creates the agent at
module scope, creates one session next to it, and serves every user from it. Everyone sees
everyone's history. It passes local testing perfectly because you are one user.

[STATELESSNESS]
Omitting `session=` is a legitimate design choice, not an oversight. Classification,
extraction and summarisation endpoints should usually be stateless — cheaper, more
predictable, trivially parallel.

[DEMO CUE]
uv run python examples/03_multi_turn.py

[PRESENTER TIP]
Run it, then comment out `session=session` on the second call and run again. The agent no
longer knows Alice. Ten seconds, perfectly clear.
""")

    # ---- 25 Context providers --------------------------------------------------------
    s = d.slide("Context providers: hooks around every run")
    d.code(s, ML, 1.85, 7.3, [
        "class UserMemoryProvider(ContextProvider):",
        "    async def before_run(self, *, agent, session, context, state):",
        "        if name := state.get(\"user_name\"):",
        "            context.extend_instructions(",
        "                self.source_id, f\"The user's name is {name}.\")",
        "",
        "    async def after_run(self, *, agent, session, context, state):",
        "        for msg in context.input_messages:",
        "            if \"my name is\" in (msg.text or \"\").lower():",
        "                state[\"user_name\"] = extract(msg.text)",
        "",
        "agent = Agent(..., context_providers=[UserMemoryProvider()])",
    ], size=12.5, title="examples/04_memory.py")
    hooks = [("before_run", "read state → inject context", GREEN),
             ("after_run", "observe the turn → write state", AMBER)]
    cy = 1.95
    for name, sub, col in hooks:
        c = d.card(s, ML + 7.7, cy, 4.1, 0.92, fill=PANEL_2, line=col)
        d.label(c, [(name, 14.5, col, True), (sub, 11.5, MUTED, False)])
        cy += 1.06
    d.text(s, ML + 7.7, 4.12, 4.1, 1.6,
           "state is namespaced per provider by source_id, so two providers never collide "
           "— and everything they write rides along in the session.", 12.5, MUTED,
           spacing=1.25)
    d.callout(s, ML, 5.66, CW, 1.02, "this is the RAG seam",
              "Retrieval is just a before_run that queries your store and calls "
              "extend_instructions. There is no separate 'RAG mode' — memory, retrieval and "
              "personalisation are all the same extension point.", CYAN)
    d.notes(s, """
[TALK TRACK]
Two hooks. before_run injects, after_run captures. That is the entire contract.

[KEY POINT — the callout]
People come in expecting a RAG API and are briefly disappointed there isn't one. Reframe
it: retrieval is a before_run hook. Personalisation is a before_run hook. Compliance
banners are a before_run hook. One seam, many uses — that is good design, not a gap.

[extend_instructions]
Note it takes source_id as the first argument. Injected context is attributed, so you can
trace which provider added what — matters when three providers are competing for a small
context window.

[BE HONEST]
The example's name extraction is a substring hack for teaching. Real implementations use a
model call, an entity extractor, or explicit user profile writes. Say so before someone in
the room says it for you.

[DEMO CUE]
uv run python examples/04_memory.py
""")

    # ---- 26 Layering providers -------------------------------------------------------
    s = d.slide("Providers stack", eyebrow="examples/05_memory_providers.py")
    d.text(s, ML, 1.78, CW, 0.34,
           "Same two hooks, different jobs. Order in the list is the order they run.",
           15, MUTED)
    layers = [
        ("InMemoryHistoryProvider", "the transcript", GREEN,
         "load_messages=True — replays history into every run"),
        ("Mem0ContextProvider", "semantic memory", VIOLET,
         "durable facts retrieved across sessions"),
        ("InMemoryHistoryProvider", "the audit copy", AMBER,
         "load_messages=False — records, never feeds the model"),
    ]
    cy = 2.32
    for name, role, col, sub in layers:
        c = d.card(s, ML, cy, 7.2, 0.98)
        d.rect(s, ML, cy, 0.055, 0.98, col)
        d.text(s, ML + 0.32, cy + 0.16, 4.4, 0.3, name, 13.5, TEXT, bold=True, font=MONO)
        d.text(s, ML + 5.1, cy + 0.17, 2.0, 0.28, role, 12, col, bold=True)
        d.text(s, ML + 0.32, cy + 0.56, 6.6, 0.3, sub, 11.5, MUTED)
        cy += 1.1
    d.code(s, ML + 7.5, 2.32, 4.3, [
        "context_providers=[",
        "    transcript,",
        "    agent_memory,",
        "    audit,   # last",
        "]",
    ], size=12)
    d.text(s, ML + 7.5, 4.3, 4.3, 1.3,
           "The audit store goes last so store_context_messages=True captures what the "
           "providers above it injected.", 12.5, MUTED, spacing=1.25)
    d.callout(s, ML, 5.7, CW, 1.0, "each provider owns a slice of session.state",
              "Namespaced by source_id, so two providers never collide. In the example the "
              "transcript holds 4 messages and the audit copy 6 — same run, different jobs.",
              CYAN)
    d.notes(s, """
[TALK TRACK]
This is the composition payoff from module 2, made concrete. Three providers, one session,
one set of hooks. Nothing here subclasses Agent.

[READ THE STACK]
The transcript is what makes the conversation multi-turn — load_messages=True replays it.
Mem0 adds durable semantic memory that outlives the session. The audit store has
load_messages=False, so it never feeds the model; it only records.

[WHY ORDER MATTERS]
Providers run in list order. The audit store is last and sets store_context_messages=True,
so it captures context the earlier providers injected. Put it first and it would record
nothing useful. That is the one detail people get wrong.

[THE NUMBERS]
When you run it, the transcript holds 4 messages and the audit copy holds 6. That gap is
the injected context, visible. Point at it.

[MEM0 IS OPTIONAL]
The example skips Mem0 unless MEM0_API_KEY is set, so it runs offline. Say that, or someone
will assume the whole slide needs a hosted account.
""")

    # ---- 27 Persistence --------------------------------------------------------------
    s = d.slide("Persistence: outliving the process",
                eyebrow="examples/06_persistent_history.py")
    a = d.card(s, ML, 1.9, 5.6, 2.15, fill=PANEL, line=CYAN)
    d.rect(s, ML, 1.9, 5.6, 0.055, CYAN)
    d.text(s, ML + 0.34, 2.14, 4.9, 0.32, "The provider persists", 17, TEXT, bold=True)
    d.text(s, ML + 0.34, 2.56, 4.9, 0.28, "FileHistoryProvider(path)", 12, CYAN, font=MONO)
    d.text(s, ML + 0.34, 2.98, 4.9, 1.3,
           "Append-only file per session. Reuse the session_id and the previous run's "
           "history is already there.\nRedis and Cosmos providers, same shape.", 12.5,
           MUTED, spacing=1.25)
    b = d.card(s, ML + 6.2, 1.9, 5.6, 2.15, fill=PANEL_2, line=GREEN)
    d.rect(s, ML + 6.2, 1.9, 5.6, 0.055, GREEN)
    d.text(s, ML + 6.54, 2.14, 4.9, 0.32, "You persist", 17, TEXT, bold=True)
    d.text(s, ML + 6.54, 2.56, 4.9, 0.28, "session.to_dict()", 12, GREEN, font=MONO)
    d.text(s, ML + 6.54, 2.98, 4.9, 1.3,
           "The session is just a dict. Serialise it into whatever store you already "
           "run, then AgentSession.from_dict() to resume.", 12.5, TEXT, spacing=1.25)
    d.code(s, ML, 4.34, CW, [
        "blob = json.dumps(session.to_dict())          # ~530 bytes",
        "restored = AgentSession.from_dict(json.loads(blob))",
        "await new_agent.run(\"What is my favourite colour?\", session=restored)",
    ], size=12.5)
    d.callout(s, ML, 5.68, CW, 0.95, "the production question this answers",
              "\"Where do sessions live when the process restarts?\" — pick one of these two "
              "before you ship, not after.", AMBER)
    d.notes(s, """
[TALK TRACK]
Module 4 opened by saying an in-memory session is fine for a CLI and wrong for a service.
This is the fix, and there are exactly two shapes.

[LEFT — the provider owns it]
FileHistoryProvider takes a directory and writes one append-only file per session. Reusing
the session_id picks up where you left off. Run the example twice and the record count goes
from 2 to 4 — that is the whole demo.

[HONEST CAVEAT]
FileHistoryProvider is marked experimental in 1.13 and warns on import. Great for local
development; for a service use Redis or Cosmos, or option two.

[RIGHT — you own it]
session.to_dict() gives you a plain dict. Serialise it wherever you already keep state.
This is usually the right answer for a web service, because you already have somewhere to
put it and you control the lifecycle and retention.

[PRESENTER TIP]
The persuasive detail is that the second half of the example builds a brand-new agent
object and hands it the restored session — and the agent still knows the favourite colour.
Same conversation, different process.
""")

    # ---- 28 Compaction ---------------------------------------------------------------
    s = d.slide("Compaction: choosing what to forget")
    d.text(s, ML, 1.78, CW, 0.34,
           "Long conversations do not fail gracefully. They fail at token limit N+1 — "
           "mid-sentence, in production.", 15, MUTED)
    strategies = [
        ("SlidingWindow", "keep the last N groups", GREEN),
        ("Truncation", "over N messages? trim to M", GREEN),
        ("ToolResultCompaction", "summarise old tool groups, keep the trace", VIOLET),
        ("SelectiveToolCall", "drop old tool groups outright", VIOLET),
        ("ContextWindow", "derive the budget from the model", CYAN),
        ("Summarization", "fold old turns into prose (needs a model)", AMBER),
    ]
    for i, (name, sub, col) in enumerate(strategies):
        col_i, row_i = i % 2, i // 2
        x = ML + col_i * 6.05
        y = 2.32 + row_i * 0.86
        d.card(s, x, y, 5.7, 0.72)
        d.rect(s, x, y, 0.055, 0.72, col)
        d.text(s, x + 0.3, y + 0.1, 3.0, 0.28, name, 13, TEXT, bold=True, font=MONO)
        d.text(s, x + 0.3, y + 0.42, 5.1, 0.26, sub, 11.5, MUTED)
    d.rect(s, ML, 4.98, CW, 0.44, PANEL_2, radius=0.1,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    d.text(s, ML, 5.06, CW, 0.32,
           "TokenBudgetComposedStrategy  —  run several, cheapest first, until the budget "
           "is met", 12.5, CYAN, bold=True, align=PP_ALIGN.CENTER)
    d.callout(s, ML, 5.66, CW, 1.05, "they are not interchangeable",
              "Dropping a tool group reclaims more tokens than summarising it — and loses "
              "the trace. Summarising costs a model call. Example 08 runs all seven side by "
              "side so you can see the trade in numbers.", AMBER)
    d.notes(s, """
[TALK TRACK]
Seven strategies, one job: decide what to forget. Do not read the list out — group it.
Two blunt ones, two tool-focused ones, one that does the budget maths for you, and one that
costs a model call.

[THE DISTINCTION THAT MATTERS]
ToolResultCompactionStrategy *replaces* old tool groups with a readable
"[Tool results: ...]" line — it reclaims the message structure and keeps a trace of what
the tools returned. SelectiveToolCallCompactionStrategy *excludes* those groups entirely.
In example 08 on the same transcript that is roughly 2,016 tokens versus 1,714. The second
saves more and tells you less. Choose deliberately.

[DO NOT OVERSELL]
Resist saying "tool results are always the biggest win" — it depends entirely on your
payload sizes. The honest framing is: measure it, which is what example 08 is for.

[COMPOSITION]
TokenBudgetComposedStrategy runs strategies in order until an included-token budget is
satisfied. Cheap structural wins first, summarisation later, hard window as a backstop.

[COST ANGLE]
For architects: compaction is a cost control as much as a correctness control. Every turn
you do not resend is money.

[WATCH OUT]
Summarisation loses detail by definition. If the agent needs exact prior values — order
IDs, amounts — keep them in structured session state, not in prose you are about to
compress.
""")

    # ---- 29 DEMO 3 -------------------------------------------------------------------
    s = d.demo(
        "Demo — sessions, providers, persistence, compaction",
        ["uv run python examples/03_multi_turn.py",
         "uv run python examples/05_memory_providers.py",
         "uv run python examples/06_persistent_history.py   # run it twice",
         "",
         "uv run python examples/08_compaction_strategies.py"],
        ["03: drop session= on turn two — the agent forgets Alice.",
         "05: transcript holds 4 messages, the audit copy 6.",
         "06: run twice — the on-disk record count grows 2 → 4.",
         "08: before/after token counts for every strategy."],
        "08 needs no model for six of its seven strategies, and 11/12/13/14 need none at "
        "all — plenty to show if the model server stalls.")
    d.notes(s, """
[SETUP]
Example 07 needs a second model server (Ollama on 11434). Decide before the session
whether you are running it; if not, present it as a code read.

[DEMO SCRIPT]
Do 03 and 04 back to back — they are fast and they build. 04 is the satisfying one: the
agent asks for the name, you give it, and two turns later it uses the name without being
prompted. Then show the printed session state so people see where it lived.

[PRESENTER TIP]
When showing 04, point out that the agent addresses Alice by name on a completely
unrelated maths question. That is before_run injecting into instructions, visibly.

[TIMING]
Five minutes for all three, less if you skip 11.
""")

    # ---- 28 SECTION 5 ----------------------------------------------------------------
    s = d.section("05", "MCP",
                  "Stop writing one-off integrations for every tool and every agent",
                  ["the N×M problem", "transports", "trust"])
    d.notes(s, """
[TALK TRACK]
Short module, high strategic value. Architects care about this one more than the syntax
modules.

[PRESENTER TIP]
If you are behind schedule, this is the module to compress to two slides — but do not skip
it. It is the most-asked-about topic in the follow-up questions.
""")

    # ---- 29 Why MCP ------------------------------------------------------------------
    s = d.slide("Model Context Protocol: N×M → N+M")
    d.text(s, ML, 1.8, 5.4, 0.3, "WITHOUT A PROTOCOL", 11, ROSE, bold=True)
    for i in range(3):
        d.card(s, ML, 2.2 + i * 0.72, 1.5, 0.55, fill=PANEL, line=STROKE)
        d.text(s, ML, 2.33 + i * 0.72, 1.5, 0.3, f"agent {i+1}", 11.5, MUTED,
               align=PP_ALIGN.CENTER)
        d.card(s, ML + 3.6, 2.2 + i * 0.72, 1.5, 0.55, fill=PANEL, line=STROKE)
        d.text(s, ML + 3.6, 2.33 + i * 0.72, 1.5, 0.3, f"tool {i+1}", 11.5, MUTED,
               align=PP_ALIGN.CENTER)
    for i in range(3):
        for j in range(3):
            d.rect(s, ML + 1.55, 2.44 + i * 0.72 + (j - 1) * 0.03, 2.0, 0.012, ROSE)
    d.text(s, ML, 4.5, 5.1, 0.3, "9 bespoke integrations · 9 places to break", 12, ROSE)
    d.text(s, ML + 6.4, 1.8, 5.4, 0.3, "WITH MCP", 11, GREEN, bold=True)
    for i in range(3):
        d.card(s, ML + 6.4, 2.2 + i * 0.72, 1.5, 0.55, fill=PANEL, line=STROKE)
        d.text(s, ML + 6.4, 2.33 + i * 0.72, 1.5, 0.3, f"agent {i+1}", 11.5, MUTED,
               align=PP_ALIGN.CENTER)
        d.card(s, ML + 10.3, 2.2 + i * 0.72, 1.5, 0.55, fill=PANEL, line=STROKE)
        d.text(s, ML + 10.3, 2.33 + i * 0.72, 1.5, 0.3, f"server {i+1}", 11.5, MUTED,
               align=PP_ALIGN.CENTER)
    hub = d.rect(s, ML + 8.35, 2.2, 1.5, 1.99, PANEL_2, line=GREEN, lw=1.4, radius=0.1,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    d.label(hub, [("MCP", 15, GREEN, True), ("one contract", 10, MUTED, False)])
    d.text(s, ML + 6.4, 4.5, 5.4, 0.3, "6 connections · one contract to test", 12, GREEN)
    d.callout(s, ML, 5.1, CW, 1.6, "why this matters beyond convenience",
              "An MCP server is a deployable, versioned, independently-owned unit. The team "
              "that owns the CRM ships the CRM tool server; your agent consumes it without "
              "a code change on either side.\n"
              "It also moves the trust boundary somewhere you can actually govern — "
              "authn, rate limits and audit live at the server, not scattered through "
              "agent code.", GREEN)
    d.notes(s, """
[TALK TRACK]
The N×M framing is the fastest way to explain MCP to engineers — it is the same argument as
ODBC, LSP, or any other protocol that collapsed a matrix into a bus.

[GO BEYOND THE DIAGRAM]
The connection-count argument is real but shallow. The deeper argument is organisational:
tool servers become owned, versioned, separately-deployed artifacts. The CRM team ships
CRM tools. You consume them. Neither side redeploys when the other changes.

[SECURITY]
Say the quiet part: an MCP server is remote code you are handing tool access to. Treat a
third-party MCP endpoint like a third-party dependency — review it, pin it, scope its
credentials, and prefer servers you or your organisation operate.

[PRESENTER TIP]
Architects will ask about auth and multi-tenancy. Honest answer: transport-level auth plus
whatever the server implements; design for the server holding the credentials, not the
agent.
""")

    # ---- 30 MCP code -----------------------------------------------------------------
    s = d.slide("Connecting a tool server", eyebrow="examples/09_mcp_docs_agent.py")
    d.code(s, ML, 1.85, 7.5, [
        "from agent_framework import Agent, MCPStreamableHTTPTool",
        "",
        "async with Agent(",
        "    client=OpenAIChatCompletionClient(),",
        "    name=\"DocsAgent\",",
        "    instructions=\"You help with Microsoft documentation questions.\",",
        "    tools=MCPStreamableHTTPTool(",
        "        name=\"Microsoft Learn MCP\",",
        "        url=\"https://learn.microsoft.com/api/mcp\",",
        "    ),",
        ") as agent:",
        "    async for chunk in agent.run(query, stream=True):",
        "        ...",
    ], size=12.5)
    d.bullets(s, ML + 7.9, 1.95, 3.9, [
        ("async with is required", "The tool owns a connection; the agent owns its lifetime."),
        ("Tools are discovered", "The server advertises its catalogue at connect time."),
        ("One line, many tools", "A server can expose dozens — you wire up none of them."),
    ], size=13.5, gap=1.0)
    d.card(s, ML + 7.9, 5.0, 3.9, 1.55, fill=PANEL_2, line=STROKE)
    d.text(s, ML + 8.18, 5.2, 3.4, 0.3, "TRANSPORTS", 10.5, CYAN, bold=True)
    d.text(s, ML + 8.18, 5.56, 3.4, 0.9,
           "Streamable HTTP for remote servers · stdio for local processes", 12, MUTED,
           spacing=1.25)
    d.notes(s, """
[TALK TRACK]
Nine lines to give an agent the entire Microsoft Learn documentation corpus as tools. No
schemas written, no endpoints wrapped.

[WHY async with]
The MCP tool holds a live connection. The context manager ties that connection's lifetime
to the agent's. Forgetting it is the number one MCP bug — you get "tool not available" or a
dangling session.

[DISCOVERY]
Emphasise that the tool list comes from the server at connect time. That is the payoff: the
server can add a tool tomorrow and your agent picks it up with no redeploy. It is also the
risk — the server can change behaviour under you.

[DEMO CUE]
uv run python examples/09_mcp_docs_agent.py

[NETWORK WARNING]
This one hits the public internet. If conference wifi is unreliable, screenshot the output
beforehand. Also note it is the slowest demo in the deck — the model reads real docs.
""")

    # ---- 31 DEMO 4 -------------------------------------------------------------------
    s = d.demo(
        "Demo — an agent with a documentation tool server",
        ["uv run python examples/09_mcp_docs_agent.py",
         "",
         "# question it answers:",
         "#  \"How do I create an Azure storage account with the az CLI?\""],
        ["No tool schemas anywhere in the file — they came from the server.",
         "The answer contains real commands, not model recall.",
         "Ask a follow-up about a service the model predates.",
         "Note the latency: remote tool + streaming generation."],
        "Requires outbound internet to learn.microsoft.com. Have a screenshot ready, "
        "and fall back to walking the nine lines of code.")
    d.notes(s, """
[DEMO SCRIPT]
Lead with the file, not the output. The persuasive thing is how little code there is.

[THE MONEY MOMENT]
Ask about something genuinely recent — newer than the model's training data. Correct answer
proves the tool did real work, which no amount of explanation achieves.

[PRESENTER TIP]
This demo is slow. Fill the gap by narrating what is happening: connect, discover tools,
model picks one, server search, results back, generation. Silence during a slow demo reads
as breakage.

[TIMING]
Three minutes, and be ready to cut it to zero if the network is bad.
""")

    # ---- 32 SECTION 6 ----------------------------------------------------------------
    s = d.section("06", "Workflows",
                  "When one agent in a loop is the wrong shape for the problem",
                  ["functional", "graph", "executors", "edges", "orchestration"])
    d.notes(s, """
[TALK TRACK]
The pivot in this module: an agent decides what to do next; a workflow is you deciding.
Both are legitimate. Knowing which you want is a design skill, and it is the thing that
separates people who ship agents from people who demo them.
""")

    # ---- 33 Agents vs workflows ------------------------------------------------------
    s = d.slide("Autonomy versus control")
    a = d.card(s, ML, 1.95, 5.6, 2.75, fill=PANEL, line=CYAN)
    d.rect(s, ML, 1.95, 5.6, 0.055, CYAN)
    d.text(s, ML + 0.34, 2.2, 4.9, 0.34, "Agent decides", 19, TEXT, bold=True)
    d.text(s, ML + 0.34, 2.66, 4.9, 1.9,
           "The model chooses the next step. Flexible, handles the unexpected, hard to "
           "predict or test.\n\nGood for: open-ended research, support triage, anything "
           "where the path is not known in advance.", 13, MUTED, spacing=1.25)
    b = d.card(s, ML + 6.2, 1.95, 5.6, 2.75, fill=PANEL_2, line=VIOLET)
    d.rect(s, ML + 6.2, 1.95, 5.6, 0.055, VIOLET)
    d.text(s, ML + 6.54, 2.2, 4.9, 0.34, "You decide", 19, TEXT, bold=True)
    d.text(s, ML + 6.54, 2.66, 4.9, 1.9,
           "The topology is code. Deterministic, testable, observable — and it will not "
           "improvise past what you wrote.\n\nGood for: pipelines, approvals, anything "
           "with an SLA or an auditor.", 13, TEXT, spacing=1.25)
    d.text(s, ML, 5.0, CW, 0.36,
           "Most real systems are both: a deterministic workflow whose nodes happen to be "
           "agents.", 17, TEXT, bold=True, font=UI_LIGHT)
    d.code(s, ML, 5.52, CW, [
        "@workflow                      # functional — plain async Python, easy to read",
        "WorkflowBuilder(...)           # graph — executors + edges, inspectable topology",
    ], size=12.5)
    d.notes(s, """
[TALK TRACK]
One axis: who chooses the next step. That is the entire distinction, and it maps directly
onto testability. A workflow you can unit test. An agent you can only evaluate.

[THE PRACTICAL ANSWER]
Nearly every production system that works is a deterministic skeleton with agentic joints.
The workflow guarantees the stages happen in order and the audit trail exists; the agents
inside each stage handle the messy language work.

[ANTI-PATTERN TO NAME]
"One big agent with twenty tools and a 2,000-word prompt telling it what order to do things
in." That is a workflow written in English instead of Python. If you find yourself
numbering steps in a system prompt, you want a workflow.

[PRESENTER TIP]
Ask for a use case from the room and classify it live. Two or three of these makes the
distinction stick far better than the slide does.
""")

    # ---- 34 Functional workflow ------------------------------------------------------
    s = d.slide("Functional workflows", eyebrow="examples/10_functional_workflow.py")
    d.code(s, ML, 1.85, 7.3, [
        "from agent_framework import Agent, workflow",
        "",
        "writer = Agent(name=\"WriterAgent\", client=client,",
        "               instructions=\"Write a short poem about the topic.\")",
        "reviewer = Agent(name=\"ReviewerAgent\", client=client,",
        "                 instructions=\"Review the poem in one sentence.\")",
        "",
        "@workflow",
        "async def poem_workflow(topic: str) -> str:",
        "    poem = (await writer.run(f\"Write about: {topic}\")).text",
        "    review = (await reviewer.run(f\"Review: {poem}\")).text",
        "    return f\"{poem}\\n\\nReview: {review}\"",
        "",
        "result = await poem_workflow.run(\"a cat learning to code\")",
        "print(result.get_outputs()[0])",
    ], size=12.5)
    d.bullets(s, ML + 7.7, 1.95, 4.1, [
        ("It reads like Python", "Because it is. Control flow is if/for/await."),
        ("Agents are just callables", "Compose them however you compose functions."),
        ("Outputs are a collection", "get_outputs() — a run can yield more than one."),
        ("Best first choice", "Reach for graphs when you need topology, not before."),
    ], size=13.5, gap=0.94)
    d.notes(s, """
[TALK TRACK]
The `@workflow` decorator buys you tracing, typed outputs and a consistent run interface
over what is otherwise ordinary async Python. The body is readable by anyone on your team
on day one.

[WHY NOT ALWAYS THIS?]
Because the topology is implicit. You cannot inspect it, visualise it, checkpoint it or
resume it mid-flight the way you can a graph. For a two-step chain that trade is obviously
correct. For a twelve-node process with retries, it is not.

[POINT AT get_outputs()]
Small API detail people trip on: the result object is not the string. A workflow run can
emit several outputs, so you index into `get_outputs()`.

[DEMO CUE]
uv run python examples/10_functional_workflow.py

[PRESENTER TIP]
Add a third agent live — an editor that rewrites based on the review. Three lines. It shows
how cheap composition is once agents are just callables.
""")

    # ---- 35 Graph workflow -----------------------------------------------------------
    s = d.slide("Graph workflows", eyebrow="examples/11_graph_workflow.py")
    d.code(s, ML, 1.85, 7.3, [
        "class UpperCase(Executor):",
        "    @handler",
        "    async def to_upper(self, text: str, ctx: WorkflowContext[str]):",
        "        await ctx.send_message(text.upper())",
        "",
        "@executor(id=\"reverse_text\")",
        "async def reverse_text(text: str, ctx: WorkflowContext[Never, str]):",
        "    await ctx.yield_output(text[::-1])",
        "",
        "upper = UpperCase(id=\"upper_case\")",
        "wf = (WorkflowBuilder(start_executor=upper)",
        "      .add_edge(upper, reverse_text)",
        "      .build())",
        "",
        "events = await wf.run(\"hello world\")   # ['DLROW OLLEH']",
    ], size=12.5)
    x = ML + 7.7
    n1 = d.card(s, x, 1.95, 4.1, 0.78, fill=PANEL, line=CYAN)
    d.label(n1, [("upper_case", 13, TEXT, True), ("Executor + @handler", 10.5, MUTED,
                                                  False)])
    d.arrow(s, x + 1.9, 2.82, 0.3, 0.4, CYAN, shape=MSO_SHAPE.DOWN_ARROW)
    n2 = d.card(s, x, 3.3, 4.1, 0.78, fill=PANEL, line=VIOLET)
    d.label(n2, [("reverse_text", 13, TEXT, True), ("@executor → yield_output", 10.5,
                                                    MUTED, False)])
    d.text(s, x, 4.3, 4.1, 1.2,
           "Nodes are executors, edges route messages. send_message passes on; "
           "yield_output ends the run with a result.", 12.5, MUTED, spacing=1.25)
    d.chip(s, x, 5.5, 3.0, 0.36, "runs with no model at all", GREEN)
    d.notes(s, """
[TALK TRACK]
Two ways to declare a node: a class with @handler methods, or a decorated async function.
Edges are declared on the builder. The topology is data, which is what makes it
inspectable and checkpointable.

[THE TWO VERBS]
`ctx.send_message(...)` forwards to the next executor. `ctx.yield_output(...)` produces a
workflow result. Confusing these is the most common first-time error — a workflow that runs
and returns nothing usually means everything called send_message and nothing yielded.

[NO MODEL REQUIRED]
Highlight the green chip. This example is pure string manipulation, which makes it the
perfect smoke test — if this runs, the install is fine and the problem is your model
server.

[TYPE PARAMETERS]
`WorkflowContext[str]` versus `WorkflowContext[Never, str]` — the first type is what this
executor sends onward, the second is what it yields as output. Never means "sends nothing
onward". Worth thirty seconds; the types are load-bearing.

[DEMO CUE]
uv run python examples/11_graph_workflow.py
""")

    # ---- 36 Executor forms -----------------------------------------------------------
    s = d.slide("Executors: four ways to declare a node",
                eyebrow="examples/12_executor_types.py")
    d.code(s, ML, 1.82, 7.15, [
        "class Ingest(Executor):                  # 1. class-based",
        "    @handler",
        "    async def from_text(self, t: str, ctx: WorkflowContext[str]):",
        "        await ctx.send_message(t)",
        "",
        "    @handler                             # 2. a second input type",
        "    async def from_lines(self, lines: list[str], ctx: ...):",
        "        await ctx.send_message(\" \".join(lines))",
        "",
        "@executor(id=\"normalise\")                 # 3. function-based",
        "async def normalise(t: str, ctx: WorkflowContext[str]):",
        "    await ctx.send_message(t.lower())",
        "",
        "@handler(input=str, output=str, workflow_output=int)   # 4. explicit",
        "async def count(self, text, ctx):        # no annotations needed",
        "    await ctx.yield_output(len(text.split()))",
    ], size=12)
    d.text(s, ML + 7.45, 1.9, 4.35, 0.3, "THE CONTEXT TYPES ARE THE CONTRACT", 10.5, CYAN,
           bold=True)
    ctxs = [("WorkflowContext", "side effects only", DIM),
            ("WorkflowContext[T]", "sends T downstream", CYAN),
            ("WorkflowContext[Never, U]", "yields U, sends nothing", GREEN),
            ("WorkflowContext[T, U]", "sends T and yields U", VIOLET)]
    cy = 2.32
    for sig, meaning, col in ctxs:
        d.rect(s, ML + 7.45, cy, 0.05, 0.62, col)
        d.text(s, ML + 7.68, cy + 0.02, 4.1, 0.26, sig, 11, TEXT, font=MONO)
        d.text(s, ML + 7.68, cy + 0.32, 4.1, 0.26, meaning, 11, MUTED)
        cy += 0.72
    d.text(s, ML + 7.45, 5.3, 4.35, 1.0,
           "Validated at build time — a mismatch fails the build, not the run.", 12.5,
           AMBER, spacing=1.25)
    d.notes(s, """
[TALK TRACK]
Four declaration forms, and one idea that matters more than all of them: the
WorkflowContext type parameters are not decoration, they are the node's contract.

[THE FOUR FORMS]
Class-based with @handler is the default. Add a second @handler and the same node accepts a
second message type — dispatch is on the message type, not the edge. The @executor
decorator gives you the same thing from a plain async function. Explicit decorator types
are the escape hatch when you cannot or will not annotate — and it is all-or-nothing, you
cannot mix explicit params with hints.

[THE CONTEXT TABLE — spend your time here]
No parameters means the node is a sink: logging, metrics, tracing. One parameter is what it
sends. WorkflowContext[Never, U] means it yields a workflow output and sends nothing — that
is your terminal node. Two parameters means it does both.

[WHY IT MATTERS]
These are validated when you call build(). Get them wrong and you get a
WorkflowValidationError before anything executes — for example, yielding an output from a
node declared WorkflowContext[str] fails with "must have output type annotations defined".
That is a good error to have seen once.

[BUILD-TIME OUTPUT DESIGNATION]
output_from and intermediate_output_from on WorkflowBuilder decide whose yields the caller
sees. There is no per-call flag and no ctx.yield_intermediate — the same yield_output call is
labelled output or intermediate purely by that build-time list.

[DEMO CUE]
uv run python examples/12_executor_types.py        (no model needed)
""")

    # ---- 37 Edge patterns ------------------------------------------------------------
    s = d.slide("Edges: six ways to wire the graph",
                eyebrow="examples/13_edge_patterns.py")
    edges = [
        ("add_chain([a, b, c])", "linear pipeline — the 90% case", CYAN),
        ("add_edge(a, b, condition=...)", "take the edge only if a predicate holds", CYAN),
        ("add_fan_out_edges(a, [b, c])", "same message to several nodes, in parallel",
         VIOLET),
        ("add_fan_in_edges([b, c], d)", "d receives a list of collected results", VIOLET),
        ("add_switch_case_edge_group(...)", "Case / Default — exactly one branch wins",
         AMBER),
        ("add_multi_selection_edge_group(...)", "pick a subset of targets at runtime",
         GREEN),
    ]
    cy = 1.88
    for sig, meaning, col in edges:
        d.card(s, ML, cy, 7.15, 0.66)
        d.rect(s, ML, cy, 0.055, 0.66, col)
        d.text(s, ML + 0.3, cy + 0.09, 4.4, 0.26, sig, 12, TEXT, bold=True, font=MONO)
        d.text(s, ML + 0.3, cy + 0.38, 6.6, 0.24, meaning, 11, MUTED)
        cy += 0.78
    d.code(s, ML + 7.45, 1.88, 4.35, [
        "wf = (WorkflowBuilder(",
        "        start_executor=triage)",
        "  .add_switch_case_edge_group(",
        "     triage, [",
        "       Case(condition=is_sev1,",
        "            target=pager),",
        "       Default(target=backlog),",
        "     ])",
        "  .build())",
    ], size=11.5)
    d.callout(s, ML + 7.45, 4.62, 4.35, 1.66, "two gotchas",
              "Conditional edges are evaluated independently — overlapping predicates run "
              "both branches.\nA switch group needs a Default.", ROSE)
    d.notes(s, """
[TALK TRACK]
Nodes do the work; edges are the design. Six primitives, and picking the right one is most
of what graph design actually is.

[WORK DOWN THE LIST]
add_chain is the one you will use most — do not be clever when a chain will do. Conditional
edges guard a branch with a predicate. Fan-out sends the same message to several nodes;
fan-in collects their results and hands the target a list, which is the bit people do not
expect — the fan-in handler signature takes list[T].

[SWITCH VS CONDITIONAL — the real distinction]
Conditional edges are independent: two overlapping predicates means both branches run. A
switch-case group picks exactly one, and requires a Default. If you want mutually exclusive
routing, use switch-case and let the framework enforce it rather than hand-writing
complementary lambdas.

[MULTI-SELECTION]
The selection function receives the message plus every candidate target id and returns the
ids you want. Good for notification fan-out where the channel set is data-driven.

[PRESENTER TIP]
Example 13 is six complete workflows in one file, all model-free, all printing their
routing decisions. It is the best file in the repo for reading aloud.

[DEMO CUE]
uv run python examples/13_edge_patterns.py        (no model needed)
""")

    # ---- 38 Workflow execution -------------------------------------------------------
    s = d.slide("Execution: supersteps, events and shared state",
                eyebrow="examples/14_workflow_execution.py")
    d.text(s, ML, 1.78, CW, 0.34,
           "Workflows advance in supersteps — each one delivers a round of messages, which "
           "is why parallel branches appear grouped.", 14.5, MUTED)
    d.code(s, ML, 2.3, 6.6, [
        "-- superstep 1",
        "   executor_invoked    parse    laptop, dock",
        "   intermediate        parse    parsed 2 items",
        "   executor_completed  parse    [['laptop','dock'], ...]",
        "-- superstep 2",
        "   executor_invoked    price    ['laptop', 'dock']",
        "   intermediate        price    priced 2 items",
        "-- superstep 3",
        "   executor_invoked    invoice  {'laptop': 60, ...}",
        "   output              invoice  invoice for 2 items ...",
    ], size=11.5, title="workflow.run(x, stream=True)")
    cards = [
        ("ctx.set_state / get_state", "scratch space every executor in the run can read",
         GREEN),
        ("get_outputs()", "the terminal answer — from output_from nodes", CYAN),
        ("get_intermediate_outputs()", "progress — from intermediate_output_from nodes",
         VIOLET),
    ]
    cy = 2.62
    for name, sub, col in cards:
        c = d.card(s, ML + 7.1, cy, 4.7, 0.92)
        d.rect(s, ML + 7.1, cy, 0.055, 0.92, col)
        d.text(s, ML + 7.4, cy + 0.14, 4.2, 0.28, name, 12.5, TEXT, bold=True, font=MONO)
        d.text(s, ML + 7.4, cy + 0.46, 4.2, 0.34, sub, 11, MUTED, spacing=1.15)
        cy += 1.04
    d.callout(s, ML, 5.72, CW, 0.95, "state is per run, not per workflow",
              "Two runs of the same workflow object do not see each other's state — so a "
              "built workflow is safe to reuse across requests.", CYAN)
    d.notes(s, """
[TALK TRACK]
Building the graph is half of it. This slide is the runtime.

[SUPERSTEPS]
This is the mental model to install. A workflow does not run node-by-node in a straight
line; it advances in supersteps, each delivering one round of messages. That is why a
fan-out's three branches all appear inside one superstep rather than interleaved. If
someone asks how parallelism is scheduled, this is the answer.

[THE EVENT STREAM]
stream=True gives you the timeline: started, superstep_started, executor_invoked,
executor_completed, output or intermediate, superstep_completed, status. Filter to
executor_invoked and output and you have a progress feed for a UI, essentially for free.
Then call get_final_response() on the stream for the aggregated result.

[SHARED STATE]
ctx.set_state / ctx.get_state — note these are synchronous in 1.13, not awaited. Use them
for facts several executors need; use messages for the thing you are actually passing
along. Do not thread a value through three handlers that do not care about it.

[THE CALLOUT MATTERS FOR ARCHITECTS]
Run state does not leak between runs, so you build the workflow once and reuse the object
per request. That is the question every architect in the room is about to ask.

[DEMO CUE]
uv run python examples/14_workflow_execution.py        (no model needed)
""")

    # ---- 39 Orchestration patterns ---------------------------------------------------
    s = d.slide("Orchestration patterns")
    pats = [
        ("Sequential", "SequentialBuilder", "A → B → C. Pipelines, review chains.", CYAN),
        ("Concurrent", "ConcurrentBuilder", "Fan out, gather, aggregate.", VIOLET),
        ("Handoff", "HandoffBuilder", "Route to the specialist, transfer control.", AMBER),
        ("Group chat", "GroupChatBuilder", "Shared transcript, managed turn-taking.", GREEN),
        ("Magentic", "MagenticBuilder", "A manager plans and delegates dynamically.", ROSE),
    ]
    for i, (name, api, desc, col) in enumerate(pats):
        col_i, row_i = i % 3, i // 3
        x = ML + col_i * 4.0
        y = 1.95 + row_i * 1.62
        d.card(s, x, y, 3.78, 1.42)
        d.rect(s, x, y, 3.78, 0.055, col)
        d.text(s, x + 0.3, y + 0.24, 3.2, 0.32, name, 16, TEXT, bold=True)
        d.text(s, x + 0.3, y + 0.62, 3.3, 0.26, api, 10.5, col, font=MONO)
        d.text(s, x + 0.3, y + 0.96, 3.2, 0.34, desc, 11.5, MUTED)
    d.card(s, ML + 8.0, 3.57, 3.78, 1.42, fill=PANEL_2, line=STROKE)
    d.text(s, ML + 8.3, 3.8, 3.2, 0.3, "PICK THE DULLEST", 10.5, CYAN, bold=True)
    d.text(s, ML + 8.3, 4.16, 3.2, 0.8,
           "Cost and unpredictability climb as you go down this list.", 12, MUTED,
           spacing=1.2)
    d.callout(s, ML, 5.5, CW, 1.0, "sequencing advice",
              "Start sequential. Move to concurrent when latency hurts. Reach for magentic "
              "only when you genuinely cannot enumerate the steps — it is the most capable "
              "and the least predictable option here.", AMBER)
    d.notes(s, """
[TALK TRACK]
Five named patterns from agent_framework.orchestrations. Most multi-agent systems are one
of the first two.

[SEQUENTIAL]
Draft → review → edit. Boring and correct.

[CONCURRENT]
Same prompt to N agents, gather all responses. Example 15 does this with a researcher, a
marketer and a compliance reviewer. Latency of the slowest, not the sum.

[HANDOFF]
Control transfers. Triage agent decides this is a billing question and hands to the billing
agent, which now owns the conversation.

[GROUP CHAT]
Agents share a transcript with managed turn-taking. Powerful, expensive, hard to bound —
every agent reads everything.

[MAGENTIC]
A manager agent plans and delegates dynamically. Impressive in demos. Be candid that
unpredictable cost and latency make it a poor default.

[PRESENTER TIP]
Land the sequencing advice. The failure mode in this space is picking the most sophisticated
pattern because it is the most interesting, then discovering your per-request cost is
twenty times what you budgeted.
""")

    # ---- 37 Concurrent code ----------------------------------------------------------
    s = d.slide("Fan-out, fan-in", eyebrow="examples/15_concurrent_orchestration.py")
    d.code(s, ML, 1.85, 7.3, [
        "from agent_framework.orchestrations import ConcurrentBuilder",
        "",
        "researcher = Agent(client=client, name=\"researcher\",",
        "                   instructions=\"Concise, factual market insights.\",",
        "                   default_options={\"max_tokens\": 500})",
        "marketer   = Agent(..., name=\"marketer\")",
        "legal      = Agent(..., name=\"legal\")",
        "",
        "workflow = ConcurrentBuilder(",
        "    participants=[researcher, marketer, legal]).build()",
        "",
        "events = await workflow.run(\"We are launching a budget e-bike.\")",
        "for msg in events.get_outputs()[0].messages:",
        "    print(msg.author_name, msg.text)",
    ], size=12.5)
    # genuinely parallel: one prompt fans out to three peers, then fans back in
    x = ML + 7.7
    src = d.card(s, x + 1.15, 2.0, 1.8, 0.6, fill=PANEL_2, line=ACCENT)
    d.label(src, [("prompt", 12.5, TEXT, True)])
    peers = [("researcher", CYAN), ("marketer", VIOLET), ("legal", AMBER)]
    for i, (n, col) in enumerate(peers):
        px = x + i * 1.41
        d.arrow(s, px + 0.53, 2.72, 0.22, 0.38, col, shape=MSO_SHAPE.DOWN_ARROW)
        c = d.card(s, px, 3.22, 1.28, 0.8, fill=PANEL, line=col)
        d.label(c, [(n, 10.5, TEXT, True)], pad=0.03)
        d.arrow(s, px + 0.53, 4.12, 0.22, 0.38, GREEN, shape=MSO_SHAPE.DOWN_ARROW)
    agg = d.card(s, x + 1.15, 4.62, 1.8, 0.6, fill=PANEL_2, line=GREEN)
    d.label(agg, [("aggregated", 12.5, TEXT, True)])
    d.text(s, x, 5.42, 4.1, 0.6, "Wall clock = the slowest participant.\n"
           "author_name tells you who said what.", 12, MUTED, spacing=1.25)
    d.notes(s, """
[TALK TRACK]
Three perspectives on one question, in parallel. The wall-clock cost is the slowest
participant, not the sum — that is the whole reason to reach for this.

[NOTE default_options]
Each agent caps max_tokens at 500 via default_options. With three agents running
concurrently against a small local model, uncapped generation will exhaust the context
budget. This is a real pattern, not a workshop artifact — bound every participant in a
fan-out.

[author_name]
The aggregated output preserves who said what. That is what makes the transcript useful
downstream, and it is why naming agents properly matters.

[DESIGN CAUTION]
Fan-out multiplies cost linearly and quality sub-linearly. Three agents is usually the
sweet spot; ten near-identical agents produce ten near-identical answers and one large
bill.

[DEMO CUE]
uv run python examples/15_concurrent_orchestration.py
""")

    # ---- 38 DEMO 5 -------------------------------------------------------------------
    s = d.demo(
        "Demo — workflows, end to end",
        ["# these four need no model at all",
         "uv run python examples/11_graph_workflow.py",
         "uv run python examples/12_executor_types.py",
         "uv run python examples/13_edge_patterns.py",
         "uv run python examples/14_workflow_execution.py",
         "",
         "uv run python examples/10_functional_workflow.py",
         "uv run python examples/15_concurrent_orchestration.py"],
        ["12: one node, two input types — dispatch is on type.",
         "13: six routing patterns, each printing its decision.",
         "14: the superstep timeline, live.",
         "10/15: agents as nodes; 15 prints author_name."],
        "Four of these are pure Python. If the model server is struggling, run 11, 12, 13 "
        "and 14, then read 10 and 15 as code with the diagrams on screen.")
    d.notes(s, """
[DEMO SCRIPT]
Order matters: the model-free ones first because they are instant, which resets confidence
after any earlier hiccup. 13 is the most quotable — six workflows in one file, each printing
its routing decision. 14 is the most instructive: the superstep timeline makes the execution
model visible in a way the diagram cannot. Then 10 for readability and 15 for the payoff.

[LIVE EDIT]
Adding an editor agent to 10 is the best live-coding moment in the deck — three lines,
obviously correct, immediately runs. Practise it so you do not fumble the f-string.

[PRESENTER TIP]
For 08, resize the terminal so all three responses fit without scrolling. Scrolling past
output while narrating loses the room.

[TIMING]
Six minutes. This is the demo worth spending on.
""")

    # ---- 39 SECTION 7 ----------------------------------------------------------------
    s = d.section("07", "Middleware",
                  "Cross-cutting concerns that do not belong in your agent code",
                  ["agent vs function", "three flavours", "short-circuit"])
    d.notes(s, """
[TALK TRACK]
If the audience knows ASP.NET Core, Express or Django middleware, say so immediately — the
mental model transfers exactly. Pipeline, call_next, short-circuit.
""")

    # ---- 40 Middleware pipeline ------------------------------------------------------
    s = d.slide("Two pipelines, two scopes")
    outer = d.rect(s, ML, 1.95, 11.0, 2.85, None, line=VIOLET, lw=1.6, radius=0.04,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    d.text(s, ML + 0.28, 2.12, 5.0, 0.3, "AGENT MIDDLEWARE  ·  once per run", 11, VIOLET,
           bold=True)
    inner = d.rect(s, ML + 0.55, 2.6, 9.9, 1.95, None, line=AMBER, lw=1.6, radius=0.05,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    d.text(s, ML + 0.85, 2.78, 6.0, 0.3, "FUNCTION MIDDLEWARE  ·  once per tool call", 11,
           AMBER, bold=True)
    core = d.card(s, ML + 1.2, 3.12, 3.1, 1.15, fill=PANEL_2, line=CYAN)
    d.label(core, [("model call", 14, TEXT, True), ("+ tool dispatch", 11, MUTED, False)])
    d.text(s, ML + 4.7, 3.32, 5.4, 0.9,
           "One agent middleware wraps the entire exchange — including every tool round "
           "trip inside it.\nOne function middleware fires per individual tool "
           "invocation.", 13, TEXT, spacing=1.3)
    uses = [("Guardrails", "block unsafe input before it reaches the model", VIOLET),
            ("Auth & tenancy", "attach identity, enforce scope per tool", AMBER),
            ("Telemetry", "spans, timings, token accounting", CYAN),
            ("Caching & retries", "short-circuit repeats, wrap flaky tools", GREEN)]
    x = ML
    for name, sub, col in uses:
        w = 2.86
        c = d.card(s, x, 5.22, w, 1.15)
        d.rect(s, x, 5.22, w, 0.05, col)
        d.text(s, x + 0.24, 5.44, w - 0.4, 0.3, name, 13.5, TEXT, bold=True)
        d.text(s, x + 0.24, 5.78, w - 0.4, 0.5, sub, 11, MUTED, spacing=1.15)
        x += w + 0.12
    d.notes(s, """
[TALK TRACK]
Two scopes, and the nesting is the whole lesson: function middleware lives inside agent
middleware. An agent middleware wrapping a run that makes six tool calls fires once; the
function middleware fires six times.

[WHY THAT MATTERS]
Put per-request concerns outside — authentication, request-level rate limits, total token
budget. Put per-tool concerns inside — argument validation, per-tool timeouts, tool-level
audit.

[THE FAMILIAR ANALOGY]
Same shape as HTTP middleware in any framework they know. `call_next()` is `next()`. Not
calling it short-circuits the pipeline.

[PRESENTER TIP]
Pose one: "you want to cap a single run at 10,000 tokens — which pipeline?" Then: "you want
to stop any single tool taking more than 5 seconds — which one?" Two questions, and the
model is locked in.
""")

    # ---- 41 Three flavours -----------------------------------------------------------
    s = d.slide("Three ways to write it")
    flavours = [
        ("Class-based", "AgentMiddleware\nFunctionMiddleware", CYAN,
         "Stateful, testable, injectable. Reach for it when the middleware has "
         "dependencies or configuration.", "examples/12"),
        ("Function-based", "async def mw(context, call_next)", VIOLET,
         "Detected from the context type hint. Lightest option for stateless "
         "concerns.", "examples/13"),
        ("Decorator", "@agent_middleware\n@function_middleware", GREEN,
         "Explicit intent, no type annotations needed. Mismatches are caught at "
         "registration.", "examples/14"),
    ]
    for i, (name, api, col, desc, ref) in enumerate(flavours):
        x = ML + i * 4.0
        d.card(s, x, 1.95, 3.78, 3.4)
        d.rect(s, x, 1.95, 3.78, 0.055, col)
        d.text(s, x + 0.3, 2.2, 3.2, 0.34, name, 17, TEXT, bold=True)
        yy = 2.66
        for ln in api.split("\n"):
            d.text(s, x + 0.3, yy, 3.3, 0.26, ln, 10.5, col, font=MONO)
            yy += 0.28
        d.text(s, x + 0.3, yy + 0.18, 3.2, 1.3, desc, 12.5, MUTED, spacing=1.25)
        d.chip(s, x + 0.3, 4.82, 1.6, 0.32, ref, DIM, size=9.5)
    d.code(s, ML, 5.62, CW, [
        "agent = Agent(..., middleware=[SecurityAgentMiddleware(),"
        " LoggingFunctionMiddleware()])",
        "# both pipelines registered in one list — the framework sorts them by type",
    ], size=12.5)
    d.notes(s, """
[TALK TRACK]
Same capability, three ergonomics. Do not spend long on the syntax — spend it on when to
pick which.

[HOW DETECTION WORKS]
The framework figures out whether something is agent or function middleware from either the
decorator or the context parameter's type annotation. Give it both and it validates they
agree. Give it neither and it raises rather than guessing — which is the right call, and
example 18's docstring spells out all four cases.

[PRACTICAL GUIDANCE]
Class-based for anything with dependencies you want to inject and mock. Function-based for
five-line stateless concerns. Decorator when you want intent legible at a glance without
importing context types.

[NOTE THE REGISTRATION]
Both kinds go in the same `middleware=` list. Nothing to separate by hand.
""")

    # ---- 42 Guardrail ----------------------------------------------------------------
    s = d.slide("Short-circuiting: a guardrail in twelve lines",
                eyebrow="examples/16_class_based_middleware.py")
    d.code(s, ML, 1.85, 7.6, [
        "class SecurityAgentMiddleware(AgentMiddleware):",
        "    async def process(self, context, call_next):",
        "        last = context.messages[-1] if context.messages else None",
        "        if last and last.text and risky(last.text):",
        "            context.result = AgentResponse(messages=[",
        "                Message(\"assistant\", [\"Request blocked.\"])])",
        "            return          # never call call_next → model never runs",
        "        await call_next()",
        "",
        "class LoggingFunctionMiddleware(FunctionMiddleware):",
        "    async def process(self, context, call_next):",
        "        start = time.time()",
        "        await call_next()",
        "        log(context.function.name, time.time() - start)",
    ], size=12.5)
    d.bullets(s, ML + 7.9, 1.95, 3.9, [
        ("Return without call_next", "The model is never invoked. No tokens spent."),
        ("Set context.result", "Otherwise the caller gets an empty response."),
        ("Time around call_next", "Before/after is where telemetry lives."),
        ("Order is the list order", "First registered is outermost."),
    ], size=13.5, gap=0.94)
    d.notes(s, """
[TALK TRACK]
Two middlewares, two shapes. The security one short-circuits; the logging one wraps.

[THE KEY LINE]
`return` without awaiting call_next. That is the short-circuit, and it means the model is
never called at all — no tokens, no latency, no chance of leaking whatever triggered the
block. Compare with post-hoc output filtering, which pays for the generation and then
throws it away.

[DO NOT FORGET context.result]
If you short-circuit without setting a result, the caller gets nothing back and it looks
like a bug. Setting a clear, user-appropriate message is part of the pattern.

[BE HONEST ABOUT THE EXAMPLE]
Substring matching for "password" and "secret" is a teaching device, not a content filter.
Real deployments use a classifier or a dedicated safety service. Say so — this audience
will notice.

[DEMO CUE]
uv run python examples/16_class_based_middleware.py
""")

    # ---- 43 DEMO 6 -------------------------------------------------------------------
    s = d.demo(
        "Demo — guardrails and timing",
        ["uv run python examples/16_class_based_middleware.py",
         "",
         "# same behaviour, different ergonomics:",
         "uv run python examples/18_decorator_middleware.py"],
        ["Normal query: security passes, tool runs, timing logged.",
         "\"What's the password…\": blocked before the model is called.",
         "Note the log ordering — outer in, inner in, inner out, outer out.",
         "18 shows the decorator form with untyped parameters."],
        "Example 17 uses FoundryChatClient and needs `az login` — use 16 and 18, which "
        "run against the local model.")
    d.notes(s, """
[IMPORTANT — repo gotcha]
examples/17 is Azure Foundry only and requires `az login`. Do not run it live unless you
have already authenticated. 16 and 18 both run against the local vLLM server.

[DEMO SCRIPT]
Run 12. The output interleaves security check, tool logging and the answer. Point at the
ordering — it is the onion, visible in stdout.

[THE BLOCKED CASE]
The second query is blocked and the tool never fires. Point out the absence of the logging
line: no model call, no tool call, no cost.

[PRESENTER TIP]
Add a print to the tool body before running so the audience can see for themselves that it
did not execute on the blocked path. Absence of output is a weak signal otherwise.

[TIMING]
Four minutes.
""")

    # ---- 44 SECTION 8 ----------------------------------------------------------------
    s = d.section("08", "Production",
                  "Structured output, hosting surfaces, observability, and the rules",
                  ["structured output", "hosting", "A2A", "telemetry", "practice"])
    d.notes(s, """
[TALK TRACK]
Energy check — this is where people start packing up. Tell them the last module is the one
their future self will thank them for, and that the best-practices slide is the one to
photograph.
""")

    # ---- 45 Structured output --------------------------------------------------------
    s = d.slide("Structured output: stop parsing prose",
                eyebrow="examples/19_response_format.py")
    d.code(s, ML, 1.85, 7.4, [
        "class OutputStruct(BaseModel):",
        "    city: str",
        "    description: str",
        "",
        "agent = OpenAIChatCompletionClient().as_agent(",
        "    name=\"CityAgent\",",
        "    instructions=\"Describe cities in a structured format.\")",
        "",
        "result = await agent.run(query, options={",
        "    \"response_format\": OutputStruct, \"max_tokens\": 500})",
        "",
        "if data := result.value:            # a parsed OutputStruct",
        "    print(data.city, data.description)",
        "else:",
        "    print(\"Failed to parse:\", result.text)",
    ], size=12.5)
    d.bullets(s, ML + 7.8, 1.95, 4.0, [
        ("result.value is typed", "A real Pydantic instance, not a dict."),
        ("Always handle the else", "Parsing can fail — especially on small models."),
        ("Works with streaming", "get_final_response() on the stream gives .value."),
        ("as_agent() shorthand", "Build an agent straight off a client."),
    ], size=13.5, gap=0.94)
    d.notes(s, """
[TALK TRACK]
This is the boundary between an agent and the rest of your system. Prose is for humans;
everything crossing a service boundary should be typed.

[result.value vs result.text]
`.value` is the parsed model; `.text` is the raw string. The walrus check in the example is
idiomatic and the else branch is mandatory — a small local model will occasionally emit
something unparseable, and you need to decide whether that is a retry or an error.

[STREAMING]
Both work together: stream tokens for perceived latency, then call
`stream.get_final_response()` for the parsed object. Example 19 shows both.

[as_agent()]
Minor but nice: `client.as_agent(...)` is shorthand for constructing an Agent around that
client. Used in example 19 and in others/01.

[PRESENTER TIP]
Frame this as the "integration slide" for anyone who has to feed an agent's output into a
downstream service. It is the answer to "how do I make this safe to consume?"
""")

    # ---- 46 Harness ------------------------------------------------------------------
    s = d.slide("Harness agents: planning built in", eyebrow="examples/20_harness_agent.py")
    d.code(s, ML, 1.85, 7.4, [
        "from agent_framework import create_harness_agent",
        "",
        "agent = create_harness_agent(OpenAIChatCompletionClient())",
        "",
        "session = agent.create_session()      # carries plan + todos + history",
        "",
        "while True:",
        "    user_input = input(\"> \")",
        "    async for chunk in agent.run(user_input, session=session,",
        "                                 stream=True):",
        "        if chunk.text:",
        "            print(chunk.text, end=\"\", flush=True)",
    ], size=12.5)
    d.card(s, ML + 7.8, 1.95, 4.0, 2.5, fill=PANEL_2, line=CYAN)
    d.text(s, ML + 8.1, 2.18, 3.4, 0.3, "WHAT THE HARNESS ADDS", 10.5, CYAN, bold=True)
    d.text(s, ML + 8.1, 2.58, 3.4, 1.7,
           "A planning loop over the base agent: it decomposes the request, tracks todos, "
           "and works through them across turns — all carried in session state.", 12.5,
           TEXT, spacing=1.25)
    d.callout(s, ML, 5.38, CW, 1.12, "when to use it — and when not to",
              "Good for open-ended, multi-step requests where you cannot enumerate the "
              "steps up front.\n"
              "Costly and hard to bound for anything you could have written as a workflow. "
              "If you can draw the flowchart, draw the flowchart.", AMBER)
    d.notes(s, """
[TALK TRACK]
One factory call gives you a planning agent. The session carries the plan and the todo
list, which is why it must be threaded through every turn.

[HOW IT RELATES TO MODULE 6]
This is the maximally autonomous end of the spectrum we drew in "autonomy versus control".
Harness agent at one end, graph workflow at the other. Same framework, opposite trade-offs.

[COST WARNING]
Planning loops burn tokens. Example 20 caps max_tokens at 200 per turn precisely because
the local model has a small window and a harness agent will happily fill it.

[HONEST TAKE]
Do not let this become the default answer. Most business processes are known processes. The
harness earns its cost when the request genuinely is open-ended — investigation, research,
debugging.

[DEMO CUE — optional]
uv run python examples/20_harness_agent.py   (interactive; only if you have time)
""")

    # ---- 47 Hosting surfaces ---------------------------------------------------------
    s = d.slide("Getting it out of your terminal")
    surfaces = [
        ("DevUI", "dev_ui/agent.py", ":8080", CYAN,
         "serve(entities=[agent]) — browser chat for local iteration."),
        ("AG-UI", "agent_ui/server.py", ":8888", VIOLET,
         "FastAPI endpoint speaking the AG-UI protocol to any front end."),
        ("A2A", "others/02_run_agent_on_a2a.py", ":9999", AMBER,
         "Publish an agent card; other agents discover and call you."),
        ("Azure Functions", "others/01_app_run.py", "serverless", GREEN,
         "AgentFunctionApp — HTTP endpoints plus durable thread state."),
    ]
    for i, (name, path, port, col, desc) in enumerate(surfaces):
        col_i, row_i = i % 2, i // 2
        x = ML + col_i * 6.1
        y = 1.95 + row_i * 1.72
        d.card(s, x, y, 5.7, 1.5)
        d.rect(s, x, y, 0.055, 1.5, col)
        d.text(s, x + 0.32, y + 0.2, 3.4, 0.32, name, 17, TEXT, bold=True)
        d.chip(s, x + 4.2, y + 0.22, 1.25, 0.3, port, col, size=9.5)
        d.text(s, x + 0.32, y + 0.62, 5.0, 0.28, path, 10.5, col, font=MONO)
        d.text(s, x + 0.32, y + 0.98, 5.0, 0.4, desc, 12, MUTED, spacing=1.15)
    d.code(s, ML, 5.5, CW, [
        "serve(entities=[agent], auto_open=True, auth_enabled=False)"
        "   # DevUI in one line",
        "add_agent_framework_fastapi_endpoint(app, agent, \"/\")"
        "         # AG-UI in one line",
    ], size=12.5)
    d.notes(s, """
[TALK TRACK]
Four surfaces, four audiences. DevUI is for you during development. AG-UI is for your front
end team. A2A is for other agents. Azure Functions is for ops.

[DEVUI]
One line and you get a browser chat with tool-call visibility. Note `auth_enabled=False` in
the repo example — that is a local-development setting and must never ship.

[AG-UI]
A standard protocol between agent backends and chat front ends, so you are not inventing a
streaming message format. The repo has both a server and a matching console client, which
makes the contract easy to demonstrate.

[AZURE FUNCTIONS]
`AgentFunctionApp` gives you HTTP endpoints plus durable thread state — which is the honest
answer to "where do sessions live in production?" for the serverless case.

[PRESENTER TIP]
DevUI is the best five-second demo in this module. Have it already running in a browser tab
so you can switch to it without waiting for a cold start.
""")

    # ---- 48 A2A ----------------------------------------------------------------------
    s = d.slide("A2A: agents as network citizens", eyebrow="others/02 · others/03")
    d.code(s, ML, 1.85, 7.0, [
        "# server — publish a card describing what this agent can do",
        "card = AgentCard(name=\"Europe Travel Agent\", version=\"1.0.0\",",
        "                 skills=[flight_skill],",
        "                 capabilities=AgentCapabilities(streaming=True))",
        "handler = DefaultRequestHandler(",
        "    agent_executor=A2AExecutor(agent, stream=True),",
        "    task_store=InMemoryTaskStore(), agent_card=card)",
        "",
        "# client — discover, then call",
        "card = await A2ACardResolver(...).get_agent_card()",
        "async with A2AAgent(name=card.name, agent_card=card,",
        "                    url=host) as remote:",
        "    response = await remote.run(\"search flights for Europe\")",
    ], size=12)
    x = ML + 7.3
    c1 = d.card(s, x, 1.95, 4.5, 0.85, fill=PANEL, line=CYAN)
    d.label(c1, [("your agent", 13, TEXT, True), ("A2AAgent client", 10.5, MUTED, False)])
    d.arrow(s, x + 2.1, 2.9, 0.28, 0.34, DIM, shape=MSO_SHAPE.DOWN_ARROW)
    c2 = d.card(s, x, 3.32, 4.5, 0.85, fill=PANEL_2, line=AMBER)
    d.label(c2, [("agent card", 13, AMBER, True),
                 ("name · skills · capabilities · url", 10.5, MUTED, False)])
    d.arrow(s, x + 2.1, 4.27, 0.28, 0.34, DIM, shape=MSO_SHAPE.DOWN_ARROW)
    c3 = d.card(s, x, 4.69, 4.5, 0.85, fill=PANEL, line=VIOLET)
    d.label(c3, [("remote agent", 13, TEXT, True), ("owned by another team", 10.5, MUTED,
                                                    False)])
    d.text(s, x, 5.72, 4.5, 0.7,
           "Discovery, not hard-coded endpoints — the card is the contract.", 12, MUTED,
           spacing=1.2)
    d.notes(s, """
[TALK TRACK]
MCP standardises agent-to-tool. A2A standardises agent-to-agent. Different problems, same
instinct: publish a contract, discover at runtime.

[THE AGENT CARD]
Name, version, skills, capabilities, supported interfaces. A remote agent advertises what
it can do and a client resolves that card before calling. This is service discovery for
agents, and it is what makes cross-team agent composition tractable.

[REPO NOTE]
others/02 serves on :9999 via Starlette; others/03 resolves the card and calls it. Two
terminals — start the server first. Worth demonstrating if time allows because the
discovery step is more convincing when you see it.

[ARCHITECT ANGLE]
The interesting question is governance: who publishes cards, where is the registry, how do
you authenticate and authorise cross-team calls. The protocol gives you the shape; the
policy is yours. Be straight that this is early and evolving.
""")

    # ---- 49 Observability ------------------------------------------------------------
    s = d.slide("You cannot operate what you cannot see")
    layers = [("Traces", "one span per run, per tool call, per model call", CYAN),
              ("Metrics", "latency, tokens, tool error rates, cost per request", VIOLET),
              ("Transcripts", "the actual messages — the only real debugging artifact",
               AMBER)]
    cy = 1.95
    for name, sub, col in layers:
        c = d.card(s, ML, cy, 6.4, 0.95)
        d.rect(s, ML, cy, 0.055, 0.95, col)
        d.text(s, ML + 0.32, cy + 0.18, 5.6, 0.3, name, 15.5, TEXT, bold=True)
        d.text(s, ML + 0.32, cy + 0.54, 5.6, 0.3, sub, 11.5, MUTED)
        cy += 1.08
    d.card(s, ML + 6.7, 1.95, 5.1, 3.11, fill=PANEL_2, line=STROKE)
    d.text(s, ML + 7.0, 2.18, 4.5, 0.3, "WHAT MAKES AGENTS DIFFERENT", 10.5, CYAN,
           bold=True)
    d.text(s, ML + 7.0, 2.6, 4.5, 2.3,
           "A failed request has one stack trace. A failed agent run has a conversation — "
           "seven model calls, four tool invocations, and a decision at step three that "
           "made everything after it wrong.\n\nWithout the transcript you are guessing.",
           13, TEXT, spacing=1.28)
    d.callout(s, ML, 5.38, CW, 1.08, "make it boring",
              "The framework emits OpenTelemetry. Send it to the same collector as the rest "
              "of your services. Agent observability should be a dashboard your SRE team "
              "already knows how to read — not a bespoke tool nobody opens.", GREEN)
    d.notes(s, """
[TALK TRACK]
Three layers, and the third is the one teams under-invest in. Traces and metrics tell you
that something was slow or expensive. Only the transcript tells you why the agent made a
bad decision.

[THE CORE DIFFERENCE]
Traditional service debugging: one request, one stack trace. Agent debugging: a
conversation with branch points. Your logging has to preserve the sequence, including the
tool arguments and results, or post-incident analysis is impossible.

[OTel]
Emphasise the standard. This is not a proprietary telemetry format — it is OpenTelemetry
into whatever collector you already run.

[PRIVACY]
Flag it before someone else does: transcripts contain user data. Redaction and retention
policy are design decisions, not afterthoughts, and middleware is a natural place to
enforce them.

[PRESENTER TIP]
Ask how many teams have a dashboard for cost per agent request. Very few hands. That is
usually the most actionable takeaway anyone gets from this slide.
""")

    # ---- 50 Local models -------------------------------------------------------------
    s = d.slide("Running against a local model", eyebrow="the workshop setup")
    d.code(s, ML, 1.85, 7.2, [
        "# vllm-run.sh",
        "vllm serve Qwen/Qwen3.5-4B \\",
        "  --max-model-len=8192 \\",
        "  --enable-prefix-caching \\",
        "  --enable-auto-tool-choice \\",
        "  --tool-call-parser qwen3_coder",
        "",
        "# .env",
        "OPENAI_API_KEY=test-key                    # placeholder, unused",
        "OPENAI_MODEL=Qwen/Qwen3.5-4B",
        "OPENAI_BASE_URL=http://localhost:8000/v1",
    ], size=12.5)
    d.bullets(s, ML + 7.6, 1.95, 4.2, [
        ("Tool calling needs flags", "--enable-auto-tool-choice and a matching parser."),
        ("8k window is the constraint", "Hence max_tokens everywhere in the examples."),
        ("Chat Completions only", "OpenAIChatCompletionClient, not OpenAIChatClient."),
        ("Same code, any backend", "Point OPENAI_BASE_URL elsewhere and nothing changes."),
    ], size=13, gap=0.9)
    d.callout(s, ML, 5.62, CW, 1.05, "why bother with a local model",
              "Free iteration, no rate limits, no data leaving the machine — and a small "
              "model is an honest test. Prompts and tool descriptions that only work on a "
              "frontier model are prompts that will surprise you later.", CYAN)
    d.notes(s, """
[TALK TRACK]
Worth a slide because it is how this whole workshop runs, and because "can I develop
against something free" is the most common practical question.

[THE FLAGS THAT MATTER]
`--enable-auto-tool-choice` plus `--tool-call-parser qwen3_coder`. Without those, tool
calling silently does not work — the model emits text that looks like a tool call and
nothing dispatches. That failure is baffling if you have not seen it before.

[THE 8K WINDOW]
This single constraint explains most of the odd-looking code in the examples: max_tokens on
individual calls, default_options on fan-out participants, an aggressive compaction budget
in example 07.

[THE ARGUMENT FOR SMALL MODELS]
Make the quality point: if your tool descriptions are vague, a frontier model papers over
it and a 4B model does not. Developing against the small model surfaces prompt weaknesses
early.

[PRESENTER TIP]
Mention prefix caching — it makes repeated system prompts much cheaper locally, which
matters when you are running the same example twenty times.
""")

    # ---- 51 Reference architecture ---------------------------------------------------
    s = d.slide("Putting it together")
    d.card(s, ML, 1.85, 2.6, 1.0, fill=PANEL_2, line=ACCENT)
    d.text(s, ML, 2.15, 2.6, 0.3, "client / front end", 12.5, TEXT, bold=True,
           align=PP_ALIGN.CENTER)
    d.text(s, ML, 2.44, 2.6, 0.3, "AG-UI · HTTP", 10.5, MUTED, align=PP_ALIGN.CENTER)
    d.arrow(s, ML + 2.68, 2.24, 0.36, 0.24, ACCENT)
    host = d.rect(s, ML + 3.25, 1.85, 5.5, 3.7, None, line=STROKE, lw=1.2, radius=0.03,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    d.text(s, ML + 3.5, 2.0, 4.0, 0.28, "YOUR SERVICE", 10, DIM, bold=True)
    mw = d.card(s, ML + 3.5, 2.36, 5.0, 0.62, fill=PANEL, line=VIOLET)
    d.label(mw, [("middleware · authn, guardrails, budgets, telemetry", 11.5, TEXT, True)])
    ag = d.card(s, ML + 3.5, 3.08, 5.0, 0.86, fill=PANEL_2, line=CYAN)
    d.label(ag, [("workflow", 13, CYAN, True),
                 ("agents as nodes · deterministic topology", 10.5, MUTED, False)])
    cp = d.card(s, ML + 3.5, 4.04, 2.4, 0.62, fill=PANEL, line=GREEN)
    d.label(cp, [("context providers", 11, TEXT, True)])
    tl = d.card(s, ML + 6.1, 4.04, 2.4, 0.62, fill=PANEL, line=AMBER)
    d.label(tl, [("tools", 11, TEXT, True)])
    d.text(s, ML + 3.5, 4.82, 5.0, 0.3, "compaction · approvals · structured output",
           10.5, DIM, align=PP_ALIGN.CENTER)
    ext = [("model provider", CYAN, 1.85), ("memory store", GREEN, 2.72),
           ("MCP servers", AMBER, 3.59), ("other agents (A2A)", VIOLET, 4.46),
           ("OTel collector", ROSE, 5.33)]
    for name, col, y in ext:
        c = d.card(s, ML + 9.5, y, 2.3, 0.72, fill=PANEL, line=col)
        d.label(c, [(name, 11.5, TEXT, True)])
        d.arrow(s, ML + 8.9, y + 0.24, 0.42, 0.24, col)
    d.text(s, ML, 6.05, 8.4, 0.4,
           "Deterministic skeleton, agentic joints, governed edges.", 16, TEXT, bold=True,
           font=UI_LIGHT)
    d.notes(s, """
[TALK TRACK]
This is the summary slide. Everything from the last eight modules, in one picture.

[READ IT OUTSIDE IN]
A front end talks to your service over AG-UI or plain HTTP. Middleware is the outermost
ring — identity, guardrails, budgets, telemetry apply before anything else happens. Inside
that, a workflow gives deterministic structure. Agents are the nodes. Context providers
and tools hang off the agents. Everything external — model, memory store, MCP servers,
other agents, telemetry — is on the far side of a governed boundary.

[THE ONE-LINER]
"Deterministic skeleton, agentic joints, governed edges." If people remember one sentence
from today, this is a good candidate.

[PRESENTER TIP]
Invite the room to place their own system on this diagram. What is missing for them is
usually middleware or observability — and that is a useful thing for them to notice out
loud before they leave.
""")

    # ---- 52 Best practices -----------------------------------------------------------
    s = d.slide("Practices that hold up")
    good = [
        ("Name things properly", "Agent names surface in traces and transcripts."),
        ("Docstrings are prompts", "Tool selection quality lives or dies here."),
        ("Bound every loop", "Max iterations, token budget, timeouts. Always."),
        ("always_require for effects", "Anything that writes, spends, or sends."),
        ("One session per conversation", "Never module-scope, never shared."),
        ("Structured output at boundaries", "Typed data leaves; prose stays inside."),
    ]
    more = [
        ("Compact before you overflow", "Design the budget before turn thirty."),
        ("Middleware for cross-cutting", "Not copy-pasted into every agent."),
        ("Workflow if you can draw it", "Autonomy is for genuinely unknown paths."),
        ("Small model in CI", "Cheap, fast, and honest about weak prompts."),
        ("Trace from day one", "Retrofitting observability into agents is misery."),
        ("Evaluate, don't eyeball", "A fixture suite beats a vibe check."),
    ]
    d.bullets(s, ML, 1.9, 5.5, good, size=14.5, gap=0.78)
    d.bullets(s, ML + 6.0, 1.9, 5.5, more, size=14.5, gap=0.78)
    d.notes(s, """
[TALK TRACK]
Tell people to photograph this one. Then actually pause for three seconds so they can.

[IF YOU ONLY EMPHASISE THREE]
"Bound every loop" — the difference between a surprising bill and a normal one.
"One session per conversation" — the cross-user data leak that passes local testing.
"Trace from day one" — the one that is genuinely painful to add later.

[ON EVALUATION]
The last item deserves a sentence even though it is not in the code samples. Agents are
stochastic; a change that looks better on the one prompt you tried is not evidence. A small
fixture suite you can re-run is the cheapest quality investment available.

[PRESENTER TIP]
Ask which of the twelve people expect to get wrong first. The answers tell you what to
follow up on in Q&A, and they are usually honest.
""")

    # ---- 53 Anti-patterns ------------------------------------------------------------
    s = d.slide("Failure modes to recognise early")
    bad = [
        ("The 2,000-word system prompt", "Numbered steps in prose. That's a workflow.",
         "→ WorkflowBuilder or @workflow"),
        ("The thirty-tool agent", "Selection accuracy collapses as the catalogue grows.",
         "→ split by domain, route with handoff"),
        ("The shared session", "One session at module scope, every user in it.",
         "→ one session per conversation"),
        ("The unbounded loop", "No iteration cap, no budget, no timeout.",
         "→ limits in middleware"),
        ("The silent truncation", "Oldest messages dropped, system prompt with them.",
         "→ a real compaction strategy"),
        ("The vibe-check rollout", "Tried five prompts, looked fine, shipped it.",
         "→ fixtures and evaluation"),
    ]
    for i, (name, why, fix) in enumerate(bad):
        col_i, row_i = i % 2, i // 2
        x = ML + col_i * 6.1
        y = 1.9 + row_i * 1.42
        d.card(s, x, y, 5.7, 1.22)
        d.rect(s, x, y, 0.055, 1.22, ROSE)
        d.text(s, x + 0.32, y + 0.16, 5.1, 0.3, name, 14.5, TEXT, bold=True)
        d.text(s, x + 0.32, y + 0.5, 5.1, 0.3, why, 11.5, MUTED)
        d.text(s, x + 0.32, y + 0.84, 5.1, 0.28, fix, 11.5, GREEN)
    d.notes(s, """
[TALK TRACK]
Same content as the previous slide inverted, because people recognise their own systems
faster in the negative form. Expect nervous laughter — that is the slide working.

[THE THIRTY-TOOL AGENT]
Worth dwelling on. Tool selection accuracy degrades as the catalogue grows, and it degrades
faster on smaller models. The fix is decomposition: several focused agents with a router,
rather than one agent that can do everything badly.

[THE SILENT TRUNCATION]
The nastiest one because it does not error. The agent just gets subtly worse as the
conversation lengthens, and by the time anyone notices, the system prompt has been falling
off for a week.

[PRESENTER TIP]
Offer amnesty: "hands up if you have shipped at least one of these." Put your own hand up
first. It makes the room honest and the Q&A far better.
""")

    # ---- 54 Workshop map -------------------------------------------------------------
    s = d.slide("The lab, end to end")
    rows = [
        ("01", "hello_agent", "agent, streaming", CYAN),
        ("02", "add_tools", "typed tools", CYAN),
        ("03", "multi_turn", "sessions", CYAN),
        ("04", "memory", "ContextProvider", CYAN),
        ("05", "memory_providers", "layered providers", CYAN),
        ("06", "persistent_history", "file + to_dict", CYAN),
        ("07", "context_compaction", "CompactionProvider", CYAN),
        ("08", "compaction_strategies", "all 7 strategies", CYAN),
        ("09", "mcp_docs_agent", "MCP over HTTP", VIOLET),
        ("10", "functional_workflow", "@workflow", VIOLET),
        ("11", "graph_workflow", "executors + edges", VIOLET),
        ("12", "executor_types", "4 executor forms", VIOLET),
        ("13", "edge_patterns", "6 wiring patterns", VIOLET),
        ("14", "workflow_execution", "supersteps, events", VIOLET),
        ("15", "concurrent_orchestration", "fan-out / fan-in", VIOLET),
        ("16", "class_based_middleware", "guardrail + timing", AMBER),
        ("17", "function_based_middleware", "async fn (Foundry)", AMBER),
        ("18", "decorator_middleware", "@agent_middleware", AMBER),
        ("19", "response_format", "structured output", GREEN),
        ("20", "harness_agent", "planning loop", GREEN),
    ]
    for i, (num, name, concept, col) in enumerate(rows):
        col_i, row_i = i // 10, i % 10
        x = ML + col_i * 6.1
        y = 1.8 + row_i * 0.475
        d.rect(s, x, y, 5.7, 0.41, PANEL if row_i % 2 == 0 else BG, radius=0.12,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        d.text(s, x + 0.2, y + 0.08, 0.5, 0.26, num, 11, col, bold=True, font=MONO)
        d.text(s, x + 0.7, y + 0.08, 3.1, 0.26, name, 11, TEXT, bold=True, font=MONO)
        d.text(s, x + 3.85, y + 0.09, 1.8, 0.24, concept, 10, MUTED)
    d.text(s, ML, 6.66, CW, 0.32,
           "plus  dev_ui/  ·  agent_ui/  ·  others/   for DevUI, AG-UI, A2A and Azure "
           "Functions hosting", 11.5, DIM)
    d.notes(s, """
[TALK TRACK]
Twenty examples in teaching order — the numbers are the order to work through them, top of
the left column to bottom of the right. Colour groups the modules: agents and memory,
workflows, middleware, production.

[POINT AT THE MODEL-FREE ONES]
11, 12, 13 and 14 need no model at all, and 08 needs one only for its summarisation
strategy. That is a third of the lab people can run before the model server finishes
warming up.

[HOMEWORK SUGGESTION]
Point people at 20 and 07 as the two worth reading afterwards — they are the most
production-relevant and the least covered by the live demos.

[REPO CAVEATS TO REPEAT]
Example 17 needs Azure Foundry and `az login`. Example 07 needs an Ollama endpoint on
11434 for the summariser. Everything else runs against the local vLLM server.

[PRESENTER TIP]
If you are running a longer session, this is the natural point to break for hands-on time.
Suggest a concrete exercise: add a second tool to 02, then make it always_require and
observe what changes.
""")

    # ---- 60 Exercises ----------------------------------------------------------------
    s = d.slide("Now you build it", eyebrow="exercises/")
    d.text(s, ML, 1.78, CW, 0.34,
           "Eight self-checking tasks. Each one fails with a hint until you finish it — "
           "five need no model.", 15, MUTED)
    ex = [
        ("01", "tool_design", "docstrings and Field descriptions as the model's API", CYAN),
        ("02", "context_provider", "before_run / after_run and provider state", CYAN),
        ("03", "compaction_budget", "composing strategies under a token budget", CYAN),
        ("04", "executor_contract", "WorkflowContext type parameters", VIOLET),
        ("05", "routing", "Case / Default switch-case routing", VIOLET),
        ("06", "fan_in", "what a fan-in target actually receives", VIOLET),
        ("07", "guardrail_middleware", "short-circuiting before the model runs", AMBER),
        ("08", "structured_output", "response_format and the failure branch", GREEN),
    ]
    for i, (num, name, what, col) in enumerate(ex):
        col_i, row_i = i // 4, i % 4
        x = ML + col_i * 6.05
        y = 2.24 + row_i * 0.72
        d.card(s, x, y, 5.7, 0.6)
        d.rect(s, x, y, 0.055, 0.6, col)
        d.text(s, x + 0.28, y + 0.05, 0.42, 0.26, num, 11, col, bold=True, font=MONO)
        d.text(s, x + 0.74, y + 0.06, 2.5, 0.26, name, 11.5, TEXT, bold=True, font=MONO)
        d.text(s, x + 0.74, y + 0.32, 4.8, 0.24, what, 10.5, MUTED)
    d.code(s, ML, 5.24, CW, [
        "uv run python exercises/01_tool_design.py      # read the FAILs, edit, repeat",
        "uv run python exercises/run_all.py --offline   # the five that need no model",
    ], size=12.5)
    d.rect(s, ML, 6.34, 0.05, 0.34, GREEN)
    d.text(s, ML + 0.24, 6.36, CW - 0.3, 0.3,
           "Worked answers in exercises/solutions/ — read them once you have something "
           "failing for the right reason.", 12, MUTED)
    d.notes(s, """
[TALK TRACK]
This is the slide to land before you let people loose. Reading code and writing code are
different skills, and only one of them survives contact with their own codebase on Monday.

[HOW THE EXERCISES WORK]
Every file has TODOs and a set of checks at the bottom. Run it and it tells you exactly
which checks fail and why. It exits non-zero until you are done, so nobody has to guess
whether they finished. Worked answers sit in exercises/solutions/.

[WHICH ONES TO ASSIGN]
If you have twenty minutes, do 01 and 04. 01 makes the "docstring is an API" lesson
mechanical — the check reads the generated JSON schema, so a vague docstring fails. 04
makes them meet a real WorkflowValidationError, which is the fastest way to internalise the
context type parameters.

[THE ONE WITH THE BEST PAYOFF]
03. It looks like "hit the token budget", but the budget is always hit — the composed
strategy falls back to dropping whole groups. The exercise makes you notice that the blunt
fallback keeps 3 messages while a tool-aware strategy keeps 8 at the same budget. That is a
real production insight and it lands much harder as a discovery than as a bullet.

[LOGISTICS]
Five of the eight need no model, so hand these out early — people can start while the model
server warms up. Point them at run_all.py --offline.
""")

    # ---- 61 Resources ----------------------------------------------------------------
    s = d.slide("Where to go next")
    res = [
        ("Documentation", "learn.microsoft.com/agent-framework",
         "Concepts, API reference, migration guidance from AutoGen and SK.", CYAN),
        ("Source & samples", "github.com/microsoft/agent-framework",
         "python/samples is far broader than this lab — read it.", VIOLET),
        ("Model Context Protocol", "modelcontextprotocol.io",
         "The spec, plus a directory of existing servers.", AMBER),
        ("This workshop", "examples/README.md",
         "Setup, the ordered lab, and stretch goals.", GREEN),
    ]
    for i, (name, url, desc, col) in enumerate(res):
        col_i, row_i = i % 2, i // 2
        x = ML + col_i * 6.1
        y = 1.95 + row_i * 1.75
        d.card(s, x, y, 5.7, 1.5)
        d.rect(s, x, y, 0.055, 1.5, col)
        d.text(s, x + 0.32, y + 0.22, 5.0, 0.32, name, 16, TEXT, bold=True)
        d.text(s, x + 0.32, y + 0.62, 5.0, 0.3, url, 12, col, font=MONO)
        d.text(s, x + 0.32, y + 1.0, 5.0, 0.4, desc, 11.5, MUTED, spacing=1.15)
    d.callout(s, ML, 5.62, CW, 1.05, "before you leave",
              "Work through exercises/ — eight self-checking tasks, five of which need no "
              "model. Then: rebuild example 10 with SequentialBuilder, or set "
              "approval_mode=\"always_require\" on a tool and watch the flow.", CYAN)
    d.notes(s, """
[TALK TRACK]
Four links. The one to push hardest is python/samples in the GitHub repo — it is
substantially broader than this lab and it tracks the framework as it moves.

[STRETCH GOALS]
These come straight from examples/README.md. The approval one is the most instructive:
flipping a single argument turns a smooth demo into a paused run that needs a human, and
seeing that mechanic is worth more than reading about it.

[PRESENTER TIP]
Have these links in the chat or on a handout before you show the slide — people photograph
URLs badly and then never follow them.
""")

    # ---- 56 Close --------------------------------------------------------------------
    s = d._blank()
    d.rect(s, 0, 0, SW, 0.09, ACCENT)
    d.rect(s, 0, 0, 4.4, 0.09, CYAN)
    d.text(s, ML, 2.0, 11.0, 0.9, "Questions", 52, TEXT, bold=True, font=UI_LIGHT)
    d.rect(s, ML, 3.12, 1.6, 0.05, ACCENT)
    d.text(s, ML, 3.5, 9.6, 0.9,
           "Deterministic skeleton. Agentic joints. Governed edges.", 26, ACCENT,
           font=UI_LIGHT, spacing=1.2)
    d.text(s, ML, 4.55, 9.6, 0.8,
           "The framework is not the interesting part — what you can now safely put in "
           "production is.", 15, MUTED, spacing=1.25)
    x = ML
    for t, col in [("uv sync", CYAN), ("./vllm-run.sh", CYAN),
                   ("uv run python examples/01_hello_agent.py", GREEN)]:
        w = 0.098 * len(t) + 0.6
        d.chip(s, x, 5.5, w, 0.4, t, col)
        x += w + 0.18
    d.rect(s, ML, 6.24, 0.05, 0.6, ACCENT)
    d.text(s, ML + 0.26, 6.22, 6.0, 0.32, "Satendra Kumar", 15, TEXT, bold=True)
    d.text(s, ML + 0.26, 6.54, 9.0, 0.28,
           "Everything in this deck is in the repo — examples/README.md is the lab guide.",
           11.5, DIM)
    d.notes(s, """
[CLOSE]
Land the one-liner one more time: deterministic skeleton, agentic joints, governed edges.
Then get out of the way and take questions.

[QUESTIONS YOU WILL GET — have answers ready]
1. "Should we migrate off AutoGen / Semantic Kernel?" — This is the forward line; migration
   guidance is in the docs; do not promise timelines you do not own.
2. "How does this compare to LangGraph / CrewAI?" — Be gracious. Differentiators worth
   naming: standards posture (MCP, A2A, AG-UI, OTel), the .NET story, and Azure hosting
   integration. Avoid trashing alternatives.
3. "Is it production ready?" — Core is 1.x; several provider packages are still beta. Say
   that plainly and let people make their own call.
4. "What does it cost to run?" — Framework is open source; you pay for models and infra.
   Point at compaction and bounded loops as the levers.
5. "Can we run it fully offline?" — Yes, that is what today's setup was. vLLM or Ollama,
   local MCP servers over stdio, no egress.

[IF THE ROOM IS QUIET]
Ask them: "what is the first thing you would build with this?" Someone always has an answer,
and it usually starts a real conversation.

[FINALLY]
Remind them the repo has all fourteen examples plus the hosting samples, and that
examples/README.md is the guided path.
""")

    WARNINGS.extend(d.warnings)
    return d.prs


def main() -> None:
    out = Path(__file__).parent / "Microsoft-Agent-Framework-Training.pptx"
    prs = build()
    for w in WARNINGS:
        print(f"  WARN  {w}")
    prs.save(out)
    print(f"Wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
